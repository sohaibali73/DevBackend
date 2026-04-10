"""
DiffEngine
==========
Detects which slides changed between two versions of a PPTX and only
re-renders the changed ones.

This is the key to "light speed" revision — when a user says "move slide 5
to position 2", only slides 2, 3, 4, 5 need to be re-previewed.  Unchanged
slides are served from cache.

How it works
------------
1. Hash each slide in the original deck using pHash (perceptual hash)
2. Hash each slide in the revised deck
3. Compare hashes → find added, removed, modified, and unchanged slides
4. Only re-render modified slides
5. Return a DiffReport + a SlideManifest for the changed slides only

Usage
-----
    from core.vision.diff_engine import DiffEngine

    engine = DiffEngine()

    report = await engine.diff(
        original_bytes=pptx_v1_bytes,
        revised_bytes=pptx_v2_bytes,
    )
    print(report.modified, report.added, report.removed)

    # Only render changed slides (fast)
    changed_manifest = await engine.render_changed(
        revised_bytes=pptx_v2_bytes,
        diff_report=report,
    )
"""

from __future__ import annotations

import asyncio
import logging
from dataclasses import dataclass, field
from typing import Any, Dict, List, Optional, Set

logger = logging.getLogger(__name__)


# =============================================================================
# Data classes
# =============================================================================

@dataclass
class SlideHash:
    """Perceptual hash + metadata for one slide."""
    index:    int
    phash:    str
    width_px: int = 0
    height_px: int = 0


@dataclass
class DiffReport:
    """Summary of differences between two PPTX versions."""
    original_count:  int
    revised_count:   int
    added:           List[int] = field(default_factory=list)    # 1-based slide indices new in revised
    removed:         List[int] = field(default_factory=list)    # 1-based indices removed from original
    modified:        List[int] = field(default_factory=list)    # 1-based indices that changed
    unchanged:       List[int] = field(default_factory=list)    # 1-based indices identical in both
    move_map:        Dict[int, int] = field(default_factory=dict)  # original_idx → revised_idx (reorders)
    hash_threshold:  float = 0.97  # similarity threshold for "unchanged"

    @property
    def has_changes(self) -> bool:
        return bool(self.added or self.removed or self.modified)

    @property
    def changed_indices(self) -> List[int]:
        """All 1-based slide indices in revised deck that need re-rendering."""
        return sorted(set(self.added + self.modified))

    def to_dict(self) -> Dict[str, Any]:
        return {
            "original_count": self.original_count,
            "revised_count":  self.revised_count,
            "added":          self.added,
            "removed":        self.removed,
            "modified":       self.modified,
            "unchanged":      self.unchanged,
            "has_changes":    self.has_changes,
            "changed_count":  len(self.changed_indices),
            "cache_hit_rate": round(len(self.unchanged) / max(self.revised_count, 1) * 100, 1),
        }


# =============================================================================
# DiffEngine
# =============================================================================

class DiffEngine:
    """
    Compares two PPTX versions at the pixel/hash level and identifies
    exactly which slides changed.

    Thread-safe. Stateless.
    """

    def __init__(self, similarity_threshold: float = 0.97):
        """
        Parameters
        ----------
        similarity_threshold : pHash similarity above which slides are
                               considered "unchanged" (0.97 = 97% similar)
        """
        self.similarity_threshold = similarity_threshold

    # ──────────────────────────────────────────────────────────────────────────
    # Public API
    # ──────────────────────────────────────────────────────────────────────────

    async def diff(
        self,
        original_bytes: bytes,
        revised_bytes:  bytes,
    ) -> DiffReport:
        """
        Compare two PPTX files and return a DiffReport.

        Parameters
        ----------
        original_bytes : raw bytes of the before-revision PPTX
        revised_bytes  : raw bytes of the after-revision PPTX

        Returns
        -------
        DiffReport with added/removed/modified/unchanged slide lists
        """
        loop = asyncio.get_event_loop()

        # Hash both decks in parallel
        orig_hashes, rev_hashes = await asyncio.gather(
            loop.run_in_executor(None, self._hash_pptx, original_bytes),
            loop.run_in_executor(None, self._hash_pptx, revised_bytes),
        )

        return self._compute_diff(orig_hashes, rev_hashes)

    async def render_changed(
        self,
        revised_bytes: bytes,
        diff_report:   DiffReport,
        dpi:           int = 150,
    ):
        """
        Render ONLY the slides identified as changed in diff_report.

        Returns a SlideManifest containing only the changed slides.
        Unchanged slides should be served from the previous job's cache.
        """
        from core.vision.slide_renderer import SlideRenderer

        changed = diff_report.changed_indices
        if not changed:
            from core.vision.slide_renderer import SlideManifest
            return SlideManifest(
                source_filename="revised",
                source_type="pptx",
                slide_count=diff_report.revised_count,
                slides=[],
                render_strategy="diff_cache_hit",
            )

        renderer = SlideRenderer(render_dpi=dpi)
        # Render only the changed slide range
        # For efficiency, render individually if changes are scattered
        if len(changed) <= 5 or (max(changed) - min(changed)) > len(changed) * 2:
            # Scattered changes — render each independently
            manifests = await asyncio.gather(*[
                renderer.render(
                    file_bytes=revised_bytes,
                    file_type="pptx",
                    filename="revised.pptx",
                    slide_range=(idx, idx),
                )
                for idx in changed
            ])
            # Merge into single manifest
            from core.vision.slide_renderer import SlideManifest
            all_slides = []
            for m in manifests:
                all_slides.extend(m.slides)
            all_slides.sort(key=lambda s: s.index)
            return SlideManifest(
                source_filename="revised",
                source_type="pptx",
                slide_count=diff_report.revised_count,
                slides=all_slides,
                render_strategy="diff_partial",
            )
        else:
            # Contiguous range — render as a block
            return await renderer.render(
                file_bytes=revised_bytes,
                file_type="pptx",
                filename="revised.pptx",
                slide_range=(min(changed), max(changed)),
            )

    async def build_full_manifest_from_cache(
        self,
        original_job_id: str,
        revised_bytes:   bytes,
        diff_report:     DiffReport,
        new_job_id:      str,
        dpi:             int = 150,
    ) -> None:
        """
        Build the complete slide set for new_job_id by:
        1. Copying cached PNGs from original_job_id for unchanged slides
        2. Rendering only changed slides from revised_bytes

        After this call, /pptx/preview/{new_job_id}/{idx} is available
        for every slide in the revised deck.
        """
        import os
        import shutil
        from pathlib import Path

        storage_root = Path(os.environ.get("STORAGE_ROOT", "/data"))
        job_store    = storage_root / "pptx_jobs"
        orig_dir     = job_store / original_job_id
        new_dir      = job_store / new_job_id
        new_dir.mkdir(parents=True, exist_ok=True)

        # ── 1. Copy unchanged slides from original cache ───────────────────────
        for idx in diff_report.unchanged:
            orig_png = orig_dir / f"slide_{idx:04d}.png"
            new_png  = new_dir  / f"slide_{idx:04d}.png"
            if orig_png.exists() and not new_png.exists():
                shutil.copy2(orig_png, new_png)

        # ── 2. Render only changed slides ─────────────────────────────────────
        changed_manifest = await self.render_changed(
            revised_bytes=revised_bytes,
            diff_report=diff_report,
            dpi=dpi,
        )
        for slide in changed_manifest.slides:
            (new_dir / f"slide_{slide.index:04d}.png").write_bytes(slide.image_bytes)

        logger.info(
            "DiffEngine: %d cached + %d rendered for job %s",
            len(diff_report.unchanged),
            len(diff_report.changed_indices),
            new_job_id,
        )

    # ──────────────────────────────────────────────────────────────────────────
    # Internal helpers
    # ──────────────────────────────────────────────────────────────────────────

    def _hash_pptx(self, pptx_bytes: bytes) -> List[SlideHash]:
        """Render all slides and compute pHashes. Blocking."""
        import asyncio
        import io

        try:
            from pptx import Presentation
            from pptx.util import Emu
            from pptx.enum.shapes import MSO_SHAPE_TYPE
            from core.vision.element_matcher import _phash
            from core.vision.slide_renderer import SlideRenderer

            prs = Presentation(io.BytesIO(pptx_bytes))
            hashes = []
            renderer = SlideRenderer(render_dpi=72)   # low DPI — just for comparison

            for i, slide in enumerate(prs.slides):
                idx = i + 1
                # Try fast embedded image extraction first
                img = renderer._try_extract_embedded_image(slide, prs)
                if img:
                    ph = _phash(img) or ""
                    hashes.append(SlideHash(index=idx, phash=ph))
                else:
                    # If no embedded image, use a content fingerprint
                    # (hash of all text + shape count)
                    fingerprint = self._text_fingerprint(slide)
                    hashes.append(SlideHash(index=idx, phash=fingerprint))

            return hashes

        except Exception as exc:
            logger.warning("DiffEngine._hash_pptx failed: %s", exc)
            return []

    @staticmethod
    def _text_fingerprint(slide) -> str:
        """Create a deterministic hash of slide text content."""
        import hashlib
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        texts.append(run.text)
        content = "|".join(texts)
        return hashlib.md5(content.encode()).hexdigest()[:16]

    def _compute_diff(
        self,
        orig_hashes: List[SlideHash],
        rev_hashes:  List[SlideHash],
    ) -> DiffReport:
        """Compare two hash lists and produce a DiffReport."""
        from core.vision.element_matcher import _hamming

        orig_count = len(orig_hashes)
        rev_count  = len(rev_hashes)

        # Build a similarity lookup matrix (orig_idx → best rev_idx + score)
        added    = []
        removed  = []
        modified = []
        unchanged = []
        move_map: Dict[int, int] = {}

        # Track which original indices have been matched
        matched_orig: Set[int] = set()

        for rev_slide in rev_hashes:
            # Find best match in original
            best_orig_idx = -1
            best_score    = 0.0

            for orig_slide in orig_hashes:
                if orig_slide.index in matched_orig:
                    continue
                if not orig_slide.phash or not rev_slide.phash:
                    continue
                # Check if it's a hex hash or a text fingerprint
                if len(orig_slide.phash) <= 16 and len(rev_slide.phash) <= 16:
                    # Text fingerprint — exact match only
                    score = 1.0 if orig_slide.phash == rev_slide.phash else 0.0
                else:
                    dist  = _hamming(orig_slide.phash, rev_slide.phash)
                    score = max(0.0, 1.0 - dist / 64.0)

                if score > best_score:
                    best_score    = score
                    best_orig_idx = orig_slide.index

            if best_score >= self.similarity_threshold and best_orig_idx >= 0:
                # This slide exists in original
                matched_orig.add(best_orig_idx)
                if best_orig_idx != rev_slide.index:
                    move_map[best_orig_idx] = rev_slide.index

                if best_score >= 0.999:
                    unchanged.append(rev_slide.index)
                else:
                    modified.append(rev_slide.index)
            else:
                # No match → new slide
                added.append(rev_slide.index)

        # Any original indices not matched → removed
        matched_orig_set = set()
        for orig_slide in orig_hashes:
            # Check if this original slide appeared in revised
            found = False
            for rev_slide in rev_hashes:
                if not orig_slide.phash or not rev_slide.phash:
                    continue
                if len(orig_slide.phash) <= 16:
                    score = 1.0 if orig_slide.phash == rev_slide.phash else 0.0
                else:
                    dist  = _hamming(orig_slide.phash, rev_slide.phash)
                    score = max(0.0, 1.0 - dist / 64.0)
                if score >= self.similarity_threshold:
                    found = True
                    break
            if not found:
                removed.append(orig_slide.index)

        return DiffReport(
            original_count=orig_count,
            revised_count=rev_count,
            added=sorted(added),
            removed=sorted(removed),
            modified=sorted(modified),
            unchanged=sorted(unchanged),
            move_map=move_map,
            hash_threshold=self.similarity_threshold,
        )
