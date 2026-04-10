"""
SlideRenderer
=============
Converts PPTX, PDF, HTML, and image files into per-slide PNG bytes.

Rendering strategies (by file type)
-------------------------------------
PPTX  → Fast path: python-pptx extracts embedded slide images directly
         (InDesign-exported static-image decks are handled this way instantly)
        Fallback: LibreOffice headless → PDF → PyMuPDF (fitz) → PNG

PDF   → PyMuPDF (fitz) — fast, no external binary needed beyond PyMuPDF itself

HTML  → Node.js / Puppeteer screenshot subprocess

Image → Pillow normalization (PNG, JPG, WEBP → normalized PNG)

Returns
-------
SlideManifest — metadata + bytes for every slide
SlideImageInfo — per-slide metadata

Usage
-----
    from core.vision.slide_renderer import SlideRenderer

    renderer = SlideRenderer()

    # From raw bytes (direct upload)
    manifest = await renderer.render(file_bytes=b"...", file_type="pptx",
                                     filename="meet_potomac.pptx")

    for slide in manifest.slides:
        print(slide.index, len(slide.image_bytes), "bytes")
"""

from __future__ import annotations

import asyncio
import base64
import io
import logging
import os
import shutil
import subprocess
import tempfile
import threading
from dataclasses import dataclass, field
from pathlib import Path
from typing import List, Optional

logger = logging.getLogger(__name__)

# ─── LibreOffice serialisation lock ──────────────────────────────────────────
# LibreOffice headless crashes on concurrent invocations sharing the same
# user profile directory.  A threading lock ensures only one conversion runs
# at a time across all async workers.
_LIBREOFFICE_LOCK = threading.Lock()


# =============================================================================
# Data classes
# =============================================================================

@dataclass
class SlideImageInfo:
    """Metadata + raw PNG bytes for a single slide."""
    index:        int            # 1-based slide number
    image_bytes:  bytes          # PNG image data
    width_px:     int = 0
    height_px:    int = 0
    aspect_ratio: str = "16:9"
    source_type:  str = ""       # "embedded" | "rendered" | "native"
    # base64-encoded version (populated lazily)
    _b64: Optional[str] = field(default=None, repr=False)

    @property
    def data_uri(self) -> str:
        """data:image/png;base64,... for embedding in HTML / JSON responses."""
        if self._b64 is None:
            self._b64 = base64.b64encode(self.image_bytes).decode()
        return f"data:image/png;base64,{self._b64}"

    @property
    def base64(self) -> str:
        if self._b64 is None:
            self._b64 = base64.b64encode(self.image_bytes).decode()
        return self._b64


@dataclass
class SlideManifest:
    """Collection of all rendered slides from a single source document."""
    source_filename: str
    source_type:     str          # "pptx" | "pdf" | "html" | "image"
    slide_count:     int
    slides:          List[SlideImageInfo] = field(default_factory=list)
    render_strategy: str = ""     # for diagnostics
    error:           Optional[str] = None

    @property
    def success(self) -> bool:
        return self.error is None and len(self.slides) > 0

    def slide(self, index: int) -> Optional[SlideImageInfo]:
        """Return slide by 1-based index."""
        for s in self.slides:
            if s.index == index:
                return s
        return None

    def slide_range(self, start: int, end: int) -> List[SlideImageInfo]:
        """Return slides[start..end] inclusive (1-based)."""
        return [s for s in self.slides if start <= s.index <= end]


# =============================================================================
# SlideRenderer
# =============================================================================

class SlideRenderer:
    """
    Thread-safe, async-compatible slide renderer.

    All heavy work is run in an executor so it does not block the event loop.
    """

    # DPI used when rasterising PDF pages via PyMuPDF or LibreOffice→PDF.
    RENDER_DPI: int = 150   # 150 DPI → ~1920×1080 for 16:9 slides, fast
    RENDER_DPI_HIGH: int = 200  # higher quality option

    def __init__(self, render_dpi: int = 150):
        self.render_dpi = render_dpi

    # ──────────────────────────────────────────────────────────────────────────
    # Public API
    # ──────────────────────────────────────────────────────────────────────────

    async def render(
        self,
        file_bytes:  bytes,
        file_type:   str,
        filename:    str = "document",
        slide_range: Optional[tuple[int, int]] = None,  # (start, end) 1-based inclusive
    ) -> SlideManifest:
        """
        Render ``file_bytes`` to per-slide PNG images.

        Parameters
        ----------
        file_bytes  : raw bytes of the uploaded file
        file_type   : "pptx" | "pdf" | "html" | "png" | "jpg" | "jpeg" | "webp"
        filename    : original filename (used for display only)
        slide_range : optional (start, end) 1-based filter; None = all slides
        """
        ftype = file_type.lower().lstrip(".")

        loop = asyncio.get_event_loop()

        if ftype in ("pptx", "ppt"):
            manifest = await loop.run_in_executor(
                None, self._render_pptx, file_bytes, filename, slide_range
            )
        elif ftype == "pdf":
            manifest = await loop.run_in_executor(
                None, self._render_pdf, file_bytes, filename, slide_range
            )
        elif ftype in ("html", "htm"):
            manifest = await loop.run_in_executor(
                None, self._render_html, file_bytes, filename
            )
        elif ftype in ("png", "jpg", "jpeg", "webp", "gif"):
            manifest = await loop.run_in_executor(
                None, self._render_image, file_bytes, filename
            )
        else:
            # Try PPTX as default
            manifest = await loop.run_in_executor(
                None, self._render_pptx, file_bytes, filename, slide_range
            )

        return manifest

    async def render_from_path(
        self,
        file_path: str,
        slide_range: Optional[tuple[int, int]] = None,
    ) -> SlideManifest:
        """Convenience wrapper — reads file from disk, then calls render()."""
        path = Path(file_path)
        if not path.exists():
            return SlideManifest(
                source_filename=path.name,
                source_type="unknown",
                slide_count=0,
                error=f"File not found: {file_path}",
            )
        file_bytes = path.read_bytes()
        return await self.render(
            file_bytes=file_bytes,
            file_type=path.suffix.lstrip("."),
            filename=path.name,
            slide_range=slide_range,
        )

    # ──────────────────────────────────────────────────────────────────────────
    # PPTX rendering
    # ──────────────────────────────────────────────────────────────────────────

    def _render_pptx(
        self,
        file_bytes:  bytes,
        filename:    str,
        slide_range: Optional[tuple[int, int]],
    ) -> SlideManifest:
        """
        Two-stage PPTX renderer.

        Stage 1 — Fast path (static-image decks):
            python-pptx iterates each slide.  If the slide contains exactly one
            picture shape that fills the entire slide, we extract the embedded
            image bytes directly.  InDesign-exported decks look exactly like
            this — no re-rendering needed, pixel-perfect fidelity at the
            resolution the image was originally saved.

        Stage 2 — Fallback (editable / mixed decks):
            LibreOffice headless converts the PPTX to PDF, then PyMuPDF
            rasterises each page at self.render_dpi.
        """
        try:
            from pptx import Presentation
            from pptx.util import Emu
        except ImportError:
            logger.error("python-pptx not installed")
            return SlideManifest(
                source_filename=filename, source_type="pptx",
                slide_count=0, error="python-pptx not installed",
            )

        try:
            prs = Presentation(io.BytesIO(file_bytes))
        except Exception as exc:
            logger.error("python-pptx failed to open %s: %s", filename, exc)
            return self._render_pptx_via_libreoffice(file_bytes, filename, slide_range)

        slide_count = len(prs.slides)
        slides: List[SlideImageInfo] = []
        strategy = "embedded"
        fallback_indices: List[int] = []   # slides that need LibreOffice

        for i, slide in enumerate(prs.slides):
            idx = i + 1  # 1-based
            if slide_range and not (slide_range[0] <= idx <= slide_range[1]):
                continue

            embedded = self._try_extract_embedded_image(slide, prs)
            if embedded:
                w, h = self._image_dimensions(embedded)
                slides.append(SlideImageInfo(
                    index=idx,
                    image_bytes=embedded,
                    width_px=w,
                    height_px=h,
                    aspect_ratio=self._aspect_ratio(w, h),
                    source_type="embedded",
                ))
            else:
                fallback_indices.append(idx)

        # ── LibreOffice fallback for non-image slides ─────────────────────────
        if fallback_indices:
            strategy = "mixed"
            fallback = self._render_pptx_via_libreoffice(
                file_bytes, filename, slide_range
            )
            if fallback.success:
                # Merge: replace any placeholder entries with rendered versions
                for fi in fallback.slides:
                    if fi.index in fallback_indices:
                        slides.append(fi)
            else:
                # Even LibreOffice failed — add blank placeholders
                for idx in fallback_indices:
                    placeholder = self._blank_slide_png(idx)
                    slides.append(SlideImageInfo(
                        index=idx, image_bytes=placeholder,
                        width_px=1920, height_px=1080,
                        source_type="placeholder",
                    ))

        slides.sort(key=lambda s: s.index)

        if not slides and slide_count == 0:
            return SlideManifest(
                source_filename=filename, source_type="pptx",
                slide_count=0, error="Presentation has no slides",
            )

        if not slides:
            # All slides filtered out by slide_range but file is valid
            return SlideManifest(
                source_filename=filename, source_type="pptx",
                slide_count=slide_count, slides=[],
                render_strategy=strategy,
            )

        return SlideManifest(
            source_filename=filename,
            source_type="pptx",
            slide_count=slide_count,
            slides=slides,
            render_strategy=strategy,
        )

    def _try_extract_embedded_image(self, slide, prs) -> Optional[bytes]:
        """
        If a slide is essentially a single full-bleed picture, extract it.
        Returns PNG bytes or None.

        Covers two common InDesign export patterns:
        1. Single <p:pic> element with an image relationship
        2. Single <p:sp> with a <a:blipFill> fill (picture placeholder)
        """
        try:
            from pptx.util import Emu
            from pptx.enum.shapes import MSO_SHAPE_TYPE

            # Collect all picture shapes (ignore logos/tiny overlays)
            slide_w = prs.slide_width
            slide_h = prs.slide_height
            big_threshold = 0.50  # shape must cover > 50% of slide area

            pictures = []
            for shape in slide.shapes:
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    shape_area = shape.width * shape.height
                    slide_area = slide_w * slide_h
                    if slide_area > 0 and (shape_area / slide_area) >= big_threshold:
                        pictures.append(shape)

            if len(pictures) == 1:
                pic = pictures[0]
                img_bytes = pic.image.blob
                # Convert to PNG if needed
                png = self._ensure_png(img_bytes, pic.image.ext)
                return png

            # Also handle blip-fill picture placeholders
            if not pictures:
                for ph in slide.placeholders:
                    try:
                        if ph.placeholder_format.type and hasattr(ph, "image"):
                            img_bytes = ph.image.blob
                            return self._ensure_png(img_bytes, ph.image.ext)
                    except Exception:
                        continue

            return None
        except Exception as exc:
            logger.debug("_try_extract_embedded_image failed: %s", exc)
            return None

    def _ensure_png(self, img_bytes: bytes, ext: str) -> bytes:
        """Convert image bytes to PNG using Pillow if needed."""
        ext_lower = (ext or "").lower().lstrip(".")
        if ext_lower == "png":
            return img_bytes
        try:
            from PIL import Image
            buf = io.BytesIO()
            img = Image.open(io.BytesIO(img_bytes)).convert("RGBA")
            img.save(buf, format="PNG")
            return buf.getvalue()
        except Exception:
            return img_bytes  # return original if conversion fails

    def _render_pptx_via_libreoffice(
        self,
        file_bytes:  bytes,
        filename:    str,
        slide_range: Optional[tuple[int, int]],
    ) -> SlideManifest:
        """
        LibreOffice headless PPTX → PDF → PyMuPDF PNG per page.
        Uses a threading lock to serialise LibreOffice calls.
        """
        temp_dir: Optional[Path] = None
        try:
            with _LIBREOFFICE_LOCK:
                temp_dir = Path(tempfile.mkdtemp(prefix="pptx_lo_"))
                in_path = temp_dir / "input.pptx"
                in_path.write_bytes(file_bytes)

                lo_profile = temp_dir / "lo_profile"
                lo_profile.mkdir()

                result = subprocess.run(
                    [
                        "soffice",
                        f"-env:UserInstallation=file://{lo_profile}",
                        "--headless",
                        "--convert-to", "pdf",
                        "--outdir", str(temp_dir),
                        str(in_path),
                    ],
                    capture_output=True,
                    timeout=120,
                )

            if result.returncode != 0:
                err = result.stderr.decode(errors="replace")
                logger.warning("LibreOffice conversion failed: %s", err)
                return SlideManifest(
                    source_filename=filename, source_type="pptx",
                    slide_count=0,
                    error=f"LibreOffice conversion failed: {err[:300]}",
                )

            pdf_path = temp_dir / "input.pdf"
            if not pdf_path.exists():
                # LibreOffice may output a differently-named file
                pdfs = list(temp_dir.glob("*.pdf"))
                if pdfs:
                    pdf_path = pdfs[0]
                else:
                    return SlideManifest(
                        source_filename=filename, source_type="pptx",
                        slide_count=0, error="LibreOffice produced no PDF",
                    )

            return self._render_pdf(pdf_path.read_bytes(), filename, slide_range)

        except FileNotFoundError:
            logger.warning("LibreOffice (soffice) not found on PATH")
            return SlideManifest(
                source_filename=filename, source_type="pptx",
                slide_count=0, error="LibreOffice not available",
            )
        except subprocess.TimeoutExpired:
            return SlideManifest(
                source_filename=filename, source_type="pptx",
                slide_count=0, error="LibreOffice timed out",
            )
        finally:
            if temp_dir and temp_dir.exists():
                shutil.rmtree(temp_dir, ignore_errors=True)

    # ──────────────────────────────────────────────────────────────────────────
    # PDF rendering
    # ──────────────────────────────────────────────────────────────────────────

    def _render_pdf(
        self,
        file_bytes:  bytes,
        filename:    str,
        slide_range: Optional[tuple[int, int]],
    ) -> SlideManifest:
        """PyMuPDF (fitz) — render PDF pages to PNG at self.render_dpi."""
        try:
            import fitz  # PyMuPDF
        except ImportError:
            return SlideManifest(
                source_filename=filename, source_type="pdf",
                slide_count=0, error="PyMuPDF not installed",
            )

        try:
            doc = fitz.open(stream=file_bytes, filetype="pdf")
        except Exception as exc:
            return SlideManifest(
                source_filename=filename, source_type="pdf",
                slide_count=0, error=f"PyMuPDF failed to open: {exc}",
            )

        page_count = len(doc)
        matrix = fitz.Matrix(self.render_dpi / 72, self.render_dpi / 72)
        slides: List[SlideImageInfo] = []

        for i in range(page_count):
            idx = i + 1
            if slide_range and not (slide_range[0] <= idx <= slide_range[1]):
                continue
            try:
                page = doc[i]
                pix = page.get_pixmap(matrix=matrix, alpha=False)
                png_bytes = pix.tobytes("png")
                slides.append(SlideImageInfo(
                    index=idx,
                    image_bytes=png_bytes,
                    width_px=pix.width,
                    height_px=pix.height,
                    aspect_ratio=self._aspect_ratio(pix.width, pix.height),
                    source_type="rendered",
                ))
            except Exception as exc:
                logger.warning("PyMuPDF failed on page %d: %s", idx, exc)
                placeholder = self._blank_slide_png(idx)
                slides.append(SlideImageInfo(
                    index=idx, image_bytes=placeholder,
                    width_px=1920, height_px=1080, source_type="placeholder",
                ))

        doc.close()

        return SlideManifest(
            source_filename=filename,
            source_type="pdf",
            slide_count=page_count,
            slides=slides,
            render_strategy="fitz",
        )

    # ──────────────────────────────────────────────────────────────────────────
    # HTML rendering
    # ──────────────────────────────────────────────────────────────────────────

    def _render_html(self, file_bytes: bytes, filename: str) -> SlideManifest:
        """Render HTML to a single PNG screenshot via Puppeteer (Node.js)."""
        temp_dir: Optional[Path] = None
        try:
            temp_dir = Path(tempfile.mkdtemp(prefix="pptx_html_"))
            html_path = temp_dir / "slide.html"
            img_path  = temp_dir / "slide.png"
            html_path.write_bytes(file_bytes)

            # Minimal Puppeteer script — no npm install needed (uses global chromium)
            puppet_script = f"""
const puppeteer = require('puppeteer');
(async () => {{
  const browser = await puppeteer.launch({{
    args: ['--no-sandbox', '--disable-setuid-sandbox']
  }});
  const page = await browser.newPage();
  await page.setViewport({{ width: 1280, height: 720 }});
  await page.goto('file://{html_path}', {{ waitUntil: 'networkidle0', timeout: 30000 }});
  await page.screenshot({{ path: '{img_path}', fullPage: false }});
  await browser.close();
}})();
"""
            script_path = temp_dir / "screenshot.js"
            script_path.write_text(puppet_script, encoding="utf-8")

            result = subprocess.run(
                ["node", str(script_path)],
                capture_output=True, timeout=60,
            )

            if result.returncode == 0 and img_path.exists():
                png_bytes = img_path.read_bytes()
                w, h = self._image_dimensions(png_bytes)
                return SlideManifest(
                    source_filename=filename, source_type="html",
                    slide_count=1,
                    slides=[SlideImageInfo(
                        index=1, image_bytes=png_bytes,
                        width_px=w, height_px=h,
                        aspect_ratio=self._aspect_ratio(w, h),
                        source_type="rendered",
                    )],
                    render_strategy="puppeteer",
                )
            else:
                err = result.stderr.decode(errors="replace")[:300]
                logger.warning("Puppeteer screenshot failed: %s", err)
                return SlideManifest(
                    source_filename=filename, source_type="html",
                    slide_count=0, error=f"Puppeteer failed: {err}",
                )
        except FileNotFoundError:
            return SlideManifest(
                source_filename=filename, source_type="html",
                slide_count=0, error="Node.js not found",
            )
        except Exception as exc:
            return SlideManifest(
                source_filename=filename, source_type="html",
                slide_count=0, error=str(exc),
            )
        finally:
            if temp_dir and temp_dir.exists():
                shutil.rmtree(temp_dir, ignore_errors=True)

    # ──────────────────────────────────────────────────────────────────────────
    # Image (single-slide) rendering
    # ──────────────────────────────────────────────────────────────────────────

    def _render_image(self, file_bytes: bytes, filename: str) -> SlideManifest:
        """Treat an uploaded image as a single slide — normalise to PNG."""
        try:
            png = self._ensure_png(file_bytes, Path(filename).suffix.lstrip("."))
            w, h = self._image_dimensions(png)
            return SlideManifest(
                source_filename=filename,
                source_type="image",
                slide_count=1,
                slides=[SlideImageInfo(
                    index=1, image_bytes=png,
                    width_px=w, height_px=h,
                    aspect_ratio=self._aspect_ratio(w, h),
                    source_type="native",
                )],
                render_strategy="pillow",
            )
        except Exception as exc:
            return SlideManifest(
                source_filename=filename, source_type="image",
                slide_count=0, error=str(exc),
            )

    # ──────────────────────────────────────────────────────────────────────────
    # Utility helpers
    # ──────────────────────────────────────────────────────────────────────────

    @staticmethod
    def _image_dimensions(png_bytes: bytes) -> tuple[int, int]:
        """Return (width, height) of PNG bytes using Pillow."""
        try:
            from PIL import Image
            img = Image.open(io.BytesIO(png_bytes))
            return img.width, img.height
        except Exception:
            return 0, 0

    @staticmethod
    def _aspect_ratio(w: int, h: int) -> str:
        if w == 0 or h == 0:
            return "16:9"
        from math import gcd
        d = gcd(w, h)
        return f"{w // d}:{h // d}"

    @staticmethod
    def _blank_slide_png(index: int) -> bytes:
        """Generate a blank grey placeholder slide PNG."""
        try:
            from PIL import Image, ImageDraw, ImageFont
            img = Image.new("RGB", (1920, 1080), color=(40, 40, 40))
            draw = ImageDraw.Draw(img)
            draw.text(
                (960, 540),
                f"Slide {index}\n(render unavailable)",
                fill=(200, 200, 200),
                anchor="mm",
            )
            buf = io.BytesIO()
            img.save(buf, format="PNG")
            return buf.getvalue()
        except Exception:
            # Fallback: minimal 1x1 grey PNG
            import struct, zlib
            def png_chunk(tag, data):
                c = zlib.crc32(tag + data) & 0xFFFFFFFF
                return struct.pack(">I", len(data)) + tag + data + struct.pack(">I", c)
            sig = b"\x89PNG\r\n\x1a\n"
            ihdr = png_chunk(b"IHDR", struct.pack(">IIBBBBB", 1, 1, 8, 2, 0, 0, 0))
            idat = png_chunk(b"IDAT", zlib.compress(b"\x00\x28\x28\x28"))
            iend = png_chunk(b"IEND", b"")
            return sig + ihdr + idat + iend

    # ──────────────────────────────────────────────────────────────────────────
    # PPTX metadata helper (slide count without rendering)
    # ──────────────────────────────────────────────────────────────────────────

    @staticmethod
    def get_slide_count(file_bytes: bytes) -> int:
        """Return slide count from a PPTX without rendering anything."""
        try:
            from pptx import Presentation
            prs = Presentation(io.BytesIO(file_bytes))
            return len(prs.slides)
        except Exception:
            return 0

    @staticmethod
    def get_pdf_page_count(file_bytes: bytes) -> int:
        """Return page count from a PDF."""
        try:
            import fitz
            doc = fitz.open(stream=file_bytes, filetype="pdf")
            count = len(doc)
            doc.close()
            return count
        except Exception:
            return 0
