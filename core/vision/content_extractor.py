"""
ContentExtractor
================
Extracts text and structured content from any document or slide image.

Sources supported
-----------------
PPTX static-image slides  → Claude Vision OCR (high accuracy on styled text)
PPTX native slides        → python-pptx text frame extraction
PDF pages                 → PyMuPDF text extraction + optional OCR fallback
DOCX                      → python-docx paragraph + table extraction
HTML                      → BeautifulSoup text extraction
Images (PNG/JPG)          → Claude Vision OCR or pytesseract fallback

Output
------
A ``DocumentContent`` object with:
  - Per-page/slide ``PageContent`` (title, bullets, body, tables, metadata)
  - Full plain-text transcript
  - Structured JSON suitable for deck planning

Usage
-----
    from core.vision.content_extractor import ContentExtractor

    extractor = ContentExtractor()

    # From PPTX (static images — uses Vision)
    content = await extractor.extract(pptx_bytes, file_type="pptx",
                                      filename="deck.pptx")
    for page in content.pages:
        print(page.title, page.bullets)

    # Get full transcript
    print(content.full_transcript)

    # Get structured JSON for deck planner
    print(content.to_dict())
"""

from __future__ import annotations

import asyncio
import base64
import io
import logging
import re
from dataclasses import dataclass, field
from typing import Any, Dict, List, Optional

logger = logging.getLogger(__name__)


# =============================================================================
# Data classes
# =============================================================================

@dataclass
class TableData:
    headers: List[str] = field(default_factory=list)
    rows:    List[List[str]] = field(default_factory=list)

    def to_markdown(self) -> str:
        if not self.headers and not self.rows:
            return ""
        parts = []
        if self.headers:
            parts.append("| " + " | ".join(self.headers) + " |")
            parts.append("| " + " | ".join(["---"] * len(self.headers)) + " |")
        for row in self.rows:
            parts.append("| " + " | ".join(str(c) for c in row) + " |")
        return "\n".join(parts)


@dataclass
class PageContent:
    """Extracted content from a single slide or document page."""
    page_index:  int                    # 1-based
    title:       str = ""
    subtitle:    str = ""
    bullets:     List[str] = field(default_factory=list)
    body_text:   str = ""              # paragraph/body text
    tables:      List[TableData] = field(default_factory=list)
    captions:    List[str] = field(default_factory=list)
    section:     str = ""              # section label or chapter
    page_type:   str = "content"       # "title" | "content" | "section_break" | "table" | "chart"
    source_type: str = "extracted"     # "ocr" | "native" | "vision"
    raw_text:    str = ""              # full OCR/extracted text dump

    @property
    def full_text(self) -> str:
        """All text on this page as a single string."""
        parts = []
        if self.title:
            parts.append(self.title)
        if self.subtitle:
            parts.append(self.subtitle)
        if self.section:
            parts.append(self.section)
        parts.extend(self.bullets)
        if self.body_text:
            parts.append(self.body_text)
        for t in self.tables:
            parts.append(t.to_markdown())
        return "\n".join(parts)

    def to_dict(self) -> Dict[str, Any]:
        return {
            "page_index": self.page_index,
            "title":      self.title,
            "subtitle":   self.subtitle,
            "section":    self.section,
            "bullets":    self.bullets,
            "body_text":  self.body_text,
            "tables":     [{"headers": t.headers, "rows": t.rows} for t in self.tables],
            "captions":   self.captions,
            "page_type":  self.page_type,
            "source_type": self.source_type,
        }


@dataclass
class DocumentContent:
    """Full extracted content from a document."""
    filename:    str
    file_type:   str
    page_count:  int
    pages:       List[PageContent] = field(default_factory=list)
    error:       Optional[str] = None

    @property
    def success(self) -> bool:
        return self.error is None and len(self.pages) > 0

    @property
    def full_transcript(self) -> str:
        """Complete text transcript of the document."""
        parts = []
        for page in self.pages:
            if page.full_text.strip():
                parts.append(f"--- Page {page.page_index} ---")
                parts.append(page.full_text.strip())
        return "\n\n".join(parts)

    @property
    def titles_list(self) -> List[str]:
        """Just the slide titles, in order."""
        return [p.title for p in self.pages if p.title]

    def to_dict(self) -> Dict[str, Any]:
        return {
            "filename":   self.filename,
            "file_type":  self.file_type,
            "page_count": self.page_count,
            "pages":      [p.to_dict() for p in self.pages],
            "error":      self.error,
        }


# =============================================================================
# ContentExtractor
# =============================================================================

class ContentExtractor:
    """
    Extracts text and structured content from documents and slide images.

    Uses a cascading strategy:
    1. Native extraction (python-pptx, python-docx, PyMuPDF) — fast and accurate
    2. Claude Vision OCR — for static-image slides or complex layouts
    3. pytesseract — fallback when Claude is unavailable
    """

    def __init__(self, use_vision: bool = True, use_tesseract: bool = True):
        self.use_vision    = use_vision
        self.use_tesseract = use_tesseract

    # ──────────────────────────────────────────────────────────────────────────
    # Public API
    # ──────────────────────────────────────────────────────────────────────────

    async def extract(
        self,
        file_bytes: bytes,
        file_type:  str,
        filename:   str = "document",
    ) -> DocumentContent:
        """
        Extract content from any file type.

        Parameters
        ----------
        file_bytes : raw file bytes
        file_type  : "pptx" | "pdf" | "docx" | "html" | "txt" | "png" | "jpg"
        filename   : display name for the file
        """
        ftype = file_type.lower().lstrip(".")

        if ftype in ("pptx", "ppt"):
            return await self._extract_pptx(file_bytes, filename)
        elif ftype == "pdf":
            return await self._extract_pdf(file_bytes, filename)
        elif ftype in ("docx", "doc"):
            return await self._extract_docx(file_bytes, filename)
        elif ftype in ("html", "htm"):
            return self._extract_html(file_bytes, filename)
        elif ftype in ("txt", "md"):
            return self._extract_text(file_bytes, filename)
        elif ftype in ("png", "jpg", "jpeg", "webp"):
            return await self._extract_image(file_bytes, filename)
        else:
            # Try PPTX as default
            return await self._extract_pptx(file_bytes, filename)

    # ──────────────────────────────────────────────────────────────────────────
    # PPTX extraction
    # ──────────────────────────────────────────────────────────────────────────

    async def _extract_pptx(self, file_bytes: bytes, filename: str) -> DocumentContent:
        """
        Two-pass PPTX extraction:
        1. python-pptx native text extraction (shapes, text frames)
        2. Claude Vision OCR for slides that have no extractable text (static images)
        """
        try:
            from pptx import Presentation
            from pptx.enum.shapes import MSO_SHAPE_TYPE, PP_PLACEHOLDER_TYPE
        except ImportError:
            return DocumentContent(filename=filename, file_type="pptx",
                                   page_count=0, error="python-pptx not installed")

        try:
            prs = Presentation(io.BytesIO(file_bytes))
        except Exception as exc:
            return DocumentContent(filename=filename, file_type="pptx",
                                   page_count=0, error=str(exc))

        pages = []
        slide_count = len(prs.slides)
        vision_needed = []  # (slide_index, image_bytes) for slides needing OCR

        for i, slide in enumerate(prs.slides):
            idx = i + 1
            page = self._extract_pptx_slide_native(slide, idx)
            pages.append(page)

            # If we got almost no text from native extraction, queue for Vision
            if self.use_vision and len(page.full_text.strip()) < 20:
                # Extract embedded image for Vision
                from core.vision.slide_renderer import SlideRenderer
                renderer = SlideRenderer()
                img = renderer._try_extract_embedded_image(slide, prs)
                if img:
                    vision_needed.append((idx, img))

        # ── Vision OCR for text-poor slides ───────────────────────────────────
        if vision_needed:
            vision_pages = await self._vision_ocr_batch(vision_needed, filename)
            for idx, vision_page in vision_pages:
                # Merge Vision results into existing page
                for j, p in enumerate(pages):
                    if p.page_index == idx:
                        pages[j] = vision_page
                        break

        return DocumentContent(
            filename=filename,
            file_type="pptx",
            page_count=slide_count,
            pages=pages,
        )

    def _extract_pptx_slide_native(self, slide, page_index: int) -> PageContent:
        """Extract text from a single PPTX slide using python-pptx."""
        from pptx.util import Emu

        title = ""
        subtitle = ""
        bullets = []
        body_parts = []
        tables = []
        section = ""
        page_type = "content"

        for shape in slide.shapes:
            if shape.has_text_frame:
                tf = shape.text_frame
                all_text = tf.text.strip()
                if not all_text:
                    continue

                shape_name = shape.name.lower()
                # Detect title shapes
                if ("title" in shape_name or
                    (hasattr(shape, "placeholder_format") and
                     shape.placeholder_format and
                     shape.placeholder_format.idx == 0)):
                    if not title:
                        title = all_text
                        page_type = "title" if page_index == 1 else "content"
                        continue

                # Extract bullets vs body
                para_texts = []
                for para in tf.paragraphs:
                    text = para.text.strip()
                    if text:
                        para_texts.append(text)

                if len(para_texts) > 1:
                    bullets.extend(para_texts)
                elif para_texts:
                    body_parts.append(para_texts[0])

            elif shape.has_table:
                tbl = shape.table
                headers = []
                rows = []
                for ri, row in enumerate(tbl.rows):
                    cells = [c.text.strip() for c in row.cells]
                    if ri == 0:
                        headers = cells
                    else:
                        rows.append(cells)
                tables.append(TableData(headers=headers, rows=rows))

        # Detect section/divider slides
        if not bullets and not body_parts and title:
            page_type = "section_break" if len(title) < 50 else "content"

        raw_text = " | ".join(filter(None, [title, subtitle] + bullets + body_parts))

        return PageContent(
            page_index=page_index,
            title=title,
            subtitle=subtitle,
            bullets=bullets,
            body_text="\n".join(body_parts),
            tables=tables,
            section=section,
            page_type=page_type,
            source_type="native",
            raw_text=raw_text,
        )

    # ──────────────────────────────────────────────────────────────────────────
    # PDF extraction
    # ──────────────────────────────────────────────────────────────────────────

    async def _extract_pdf(self, file_bytes: bytes, filename: str) -> DocumentContent:
        """Extract text from PDF using PyMuPDF text extraction."""
        try:
            import fitz
        except ImportError:
            return DocumentContent(filename=filename, file_type="pdf",
                                   page_count=0, error="PyMuPDF not installed")

        try:
            doc = fitz.open(stream=file_bytes, filetype="pdf")
        except Exception as exc:
            return DocumentContent(filename=filename, file_type="pdf",
                                   page_count=0, error=str(exc))

        pages = []
        for i in range(len(doc)):
            page = doc[i]
            blocks = page.get_text("blocks")  # list of (x0,y0,x1,y1,text,...)
            texts  = [b[4].strip() for b in blocks if b[4].strip()]

            title = texts[0] if texts else ""
            body  = "\n".join(texts[1:]) if len(texts) > 1 else ""
            bullets = []

            # Heuristic: lines starting with bullet characters → bullets
            body_lines = []
            for t in texts[1:]:
                for line in t.split("\n"):
                    line = line.strip()
                    if line.startswith(("•", "·", "-", "*", "–", "—")) or \
                       (len(line) > 2 and line[0].isdigit() and line[1] in ".):"):
                        bullets.append(line.lstrip("•·-*–— 0123456789.)").strip())
                    elif line:
                        body_lines.append(line)

            pages.append(PageContent(
                page_index=i + 1,
                title=title,
                bullets=bullets,
                body_text="\n".join(body_lines),
                page_type="content",
                source_type="native",
                raw_text="\n".join(texts),
            ))

        doc.close()

        return DocumentContent(
            filename=filename,
            file_type="pdf",
            page_count=len(pages),
            pages=pages,
        )

    # ──────────────────────────────────────────────────────────────────────────
    # DOCX extraction
    # ──────────────────────────────────────────────────────────────────────────

    async def _extract_docx(self, file_bytes: bytes, filename: str) -> DocumentContent:
        """Extract text from DOCX using python-docx."""
        try:
            from docx import Document as DocxDocument
        except ImportError:
            return DocumentContent(filename=filename, file_type="docx",
                                   page_count=0, error="python-docx not installed")

        try:
            doc = DocxDocument(io.BytesIO(file_bytes))
        except Exception as exc:
            return DocumentContent(filename=filename, file_type="docx",
                                   page_count=0, error=str(exc))

        # Group paragraphs into "pages" (by heading level)
        pages = []
        current_title = ""
        current_bullets = []
        current_body = []
        page_idx = 1

        for para in doc.paragraphs:
            text = para.text.strip()
            if not text:
                continue
            style = para.style.name.lower() if para.style else ""

            if "heading 1" in style or "title" in style:
                # Save current page
                if current_title or current_bullets or current_body:
                    pages.append(PageContent(
                        page_index=page_idx,
                        title=current_title,
                        bullets=current_bullets,
                        body_text="\n".join(current_body),
                        page_type="content",
                        source_type="native",
                    ))
                    page_idx += 1
                current_title  = text
                current_bullets = []
                current_body   = []
            elif "list" in style or "bullet" in style:
                current_bullets.append(text)
            else:
                current_body.append(text)

        # Final page
        if current_title or current_bullets or current_body:
            pages.append(PageContent(
                page_index=page_idx,
                title=current_title,
                bullets=current_bullets,
                body_text="\n".join(current_body),
                page_type="content",
                source_type="native",
            ))

        # Also extract tables
        for tbl in doc.tables:
            rows_data = []
            for row in tbl.rows:
                rows_data.append([c.text.strip() for c in row.cells])
            if rows_data:
                pages.append(PageContent(
                    page_index=page_idx,
                    tables=[TableData(
                        headers=rows_data[0] if rows_data else [],
                        rows=rows_data[1:] if len(rows_data) > 1 else [],
                    )],
                    page_type="table",
                    source_type="native",
                ))
                page_idx += 1

        return DocumentContent(
            filename=filename,
            file_type="docx",
            page_count=len(pages),
            pages=pages,
        )

    # ──────────────────────────────────────────────────────────────────────────
    # HTML / text extraction
    # ──────────────────────────────────────────────────────────────────────────

    def _extract_html(self, file_bytes: bytes, filename: str) -> DocumentContent:
        """Extract text from HTML using BeautifulSoup or regex fallback."""
        text = file_bytes.decode(errors="replace")

        try:
            from bs4 import BeautifulSoup
            soup = BeautifulSoup(text, "html.parser")
            # Remove script/style
            for tag in soup(["script", "style"]):
                tag.decompose()
            clean = soup.get_text(separator="\n", strip=True)
        except ImportError:
            # Simple regex fallback
            clean = re.sub(r"<[^>]+>", " ", text)
            clean = re.sub(r"\s+", " ", clean).strip()

        lines = [l.strip() for l in clean.split("\n") if l.strip()]
        title = lines[0] if lines else ""
        body  = "\n".join(lines[1:]) if len(lines) > 1 else ""

        return DocumentContent(
            filename=filename,
            file_type="html",
            page_count=1,
            pages=[PageContent(
                page_index=1,
                title=title,
                body_text=body,
                source_type="native",
                raw_text=clean,
            )],
        )

    def _extract_text(self, file_bytes: bytes, filename: str) -> DocumentContent:
        """Extract text from plain text / markdown."""
        text = file_bytes.decode(errors="replace")
        lines = [l.strip() for l in text.split("\n") if l.strip()]
        title = ""
        bullets = []
        body_parts = []

        for line in lines:
            if line.startswith("# "):
                title = line[2:]
            elif line.startswith(("- ", "* ", "• ")):
                bullets.append(line[2:].strip())
            elif line.startswith("## "):
                body_parts.append(f"\n{line[3:]}")
            else:
                body_parts.append(line)

        return DocumentContent(
            filename=filename,
            file_type="txt",
            page_count=1,
            pages=[PageContent(
                page_index=1,
                title=title,
                bullets=bullets,
                body_text="\n".join(body_parts),
                source_type="native",
                raw_text=text,
            )],
        )

    # ──────────────────────────────────────────────────────────────────────────
    # Image OCR (single image)
    # ──────────────────────────────────────────────────────────────────────────

    async def _extract_image(self, image_bytes: bytes, filename: str) -> DocumentContent:
        """Extract text from an image using Claude Vision or pytesseract."""
        if self.use_vision:
            pages = await self._vision_ocr_batch([(1, image_bytes)], filename)
            if pages:
                _, page = pages[0]
                return DocumentContent(
                    filename=filename, file_type="image",
                    page_count=1, pages=[page],
                )

        if self.use_tesseract:
            text = await asyncio.get_event_loop().run_in_executor(
                None, self._tesseract_ocr, image_bytes
            )
            return DocumentContent(
                filename=filename, file_type="image",
                page_count=1,
                pages=[PageContent(
                    page_index=1, body_text=text,
                    source_type="ocr", raw_text=text,
                )],
            )

        return DocumentContent(
            filename=filename, file_type="image",
            page_count=0, error="No OCR backend available",
        )

    # ──────────────────────────────────────────────────────────────────────────
    # Vision OCR backend
    # ──────────────────────────────────────────────────────────────────────────

    async def _vision_ocr_batch(
        self,
        slides: List[tuple],     # [(slide_index, image_bytes), ...]
        context: str = "",
    ) -> List[tuple]:             # [(slide_index, PageContent), ...]
        """Send slide images to Claude Vision and extract structured text."""
        import os
        api_key = os.environ.get("ANTHROPIC_API_KEY", "")
        if not api_key:
            return []

        results = []
        sem = asyncio.Semaphore(3)

        async def _extract_one(idx: int, img_bytes: bytes) -> tuple:
            async with sem:
                try:
                    import anthropic
                    b64 = base64.b64encode(img_bytes).decode()
                    client = anthropic.AsyncAnthropic(api_key=api_key)
                    msg = await client.messages.create(
                        model="claude-sonnet-4-5",
                        max_tokens=1024,
                        system=(
                            "You are a precise text extractor. Extract ALL text from this "
                            "presentation slide image. Return ONLY a JSON object with these fields:\n"
                            '{ "title": "<slide title>", "subtitle": "<subtitle if any>", '
                            '"section": "<section label if any>", '
                            '"bullets": ["<bullet 1>", ...], '
                            '"body_text": "<any paragraph text>", '
                            '"page_type": "<title|content|section_break|table|chart>" }\n'
                            "Return ONLY valid JSON, no prose."
                        ),
                        messages=[{
                            "role": "user",
                            "content": [
                                {
                                    "type": "image",
                                    "source": {"type": "base64", "media_type": "image/png", "data": b64},
                                },
                                {"type": "text", "text": f"Extract all text from slide {idx}. JSON only."},
                            ],
                        }],
                    )
                    raw = msg.content[0].text.strip()
                    raw = re.sub(r"^```(?:json)?", "", raw).strip()
                    raw = re.sub(r"```$", "", raw).strip()
                    import json
                    data = json.loads(raw)
                    page = PageContent(
                        page_index=idx,
                        title=data.get("title", ""),
                        subtitle=data.get("subtitle", ""),
                        section=data.get("section", ""),
                        bullets=data.get("bullets") or [],
                        body_text=data.get("body_text", ""),
                        page_type=data.get("page_type", "content"),
                        source_type="vision",
                    )
                    return (idx, page)
                except Exception as exc:
                    logger.warning("Vision OCR failed for slide %d: %s", idx, exc)
                    return (idx, PageContent(page_index=idx, source_type="vision",
                                             error=str(exc) if hasattr(PageContent, "error") else ""))

        tasks = [_extract_one(idx, img) for idx, img in slides]
        raw_results = await asyncio.gather(*tasks, return_exceptions=True)
        for r in raw_results:
            if isinstance(r, tuple):
                results.append(r)

        return results

    # ──────────────────────────────────────────────────────────────────────────
    # Tesseract fallback
    # ──────────────────────────────────────────────────────────────────────────

    @staticmethod
    def _tesseract_ocr(image_bytes: bytes) -> str:
        try:
            import pytesseract
            from PIL import Image
            img = Image.open(io.BytesIO(image_bytes))
            return pytesseract.image_to_string(img)
        except Exception as exc:
            logger.warning("tesseract OCR failed: %s", exc)
            return ""
