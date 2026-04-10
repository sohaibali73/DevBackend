"""
PPTX Reviser
============
Apply targeted revisions to any existing .pptx file without regenerating it.
Uses python-pptx (pure Python, no subprocess, instant execution).

The single biggest time-saver: junior IB analysts spend 80% of their time
updating numbers in existing decks. This tool does an entire data refresh
in < 500ms.

Operations
----------
find_replace     → {"type": "find_replace", "find": "Q1 2025", "replace": "Q2 2025"}
update_slide     → {"type": "update_slide", "slide_index": 3, "slide": {generate_pptx slide spec}}
append_slides    → {"type": "append_slides", "slides": [...]}
delete_slide     → {"type": "delete_slide", "slide_index": 5}
reorder_slides   → {"type": "reorder_slides", "order": [0,2,1,3,4]}
update_table     → {"type": "update_table", "slide_index": 8, "row": 2, "col": 3, "value": "14.2x"}

Usage
-----
    from core.sandbox.pptx_reviser import PptxReviser

    reviser = PptxReviser()
    result  = reviser.revise(pptx_bytes, revisions, output_filename="Updated.pptx")
    if result.success:
        # result.data contains the updated .pptx bytes
"""

from __future__ import annotations

import copy
import io
import logging
import time
from dataclasses import dataclass
from typing import Any, Dict, List, Optional

logger = logging.getLogger(__name__)


@dataclass
class ReviseResult:
    success: bool
    error: Optional[str] = None
    exec_time_ms: float = 0.0
    data: Optional[bytes] = None
    filename: Optional[str] = None
    operations_applied: int = 0
    replacements_made: int = 0


class PptxReviser:
    """
    Apply targeted edits to an existing .pptx without regenerating it.

    Thread-safe. Stateless.
    """

    def revise(
        self,
        pptx_bytes: bytes,
        revisions: List[Dict[str, Any]],
        output_filename: Optional[str] = None,
    ) -> ReviseResult:
        start = time.time()
        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt
            from lxml import etree

            prs = Presentation(io.BytesIO(pptx_bytes))

            operations_applied = 0
            replacements_made  = 0

            for rev in revisions:
                op = rev.get("type", "")

                # ── find_replace ───────────────────────────────────────────
                if op == "find_replace":
                    find    = str(rev.get("find",    ""))
                    replace = str(rev.get("replace", ""))
                    if not find:
                        continue
                    count = self._find_replace_all(prs, find, replace)
                    replacements_made  += count
                    operations_applied += 1
                    logger.debug("find_replace %r → %r: %d replacements", find, replace, count)

                # ── delete_slide ───────────────────────────────────────────
                elif op == "delete_slide":
                    idx = rev.get("slide_index", -1)
                    if 0 <= idx < len(prs.slides):
                        xml_slides = prs.slides._sldIdLst
                        slide_elem = xml_slides[idx]
                        xml_slides.remove(slide_elem)
                        operations_applied += 1
                        logger.debug("delete_slide: index=%d", idx)

                # ── reorder_slides ─────────────────────────────────────────
                elif op == "reorder_slides":
                    order = rev.get("order", [])
                    if order and len(order) == len(prs.slides):
                        xml_slides = prs.slides._sldIdLst
                        elems = list(xml_slides)
                        for elem in elems:
                            xml_slides.remove(elem)
                        for i in order:
                            if 0 <= i < len(elems):
                                xml_slides.append(elems[i])
                        operations_applied += 1
                        logger.debug("reorder_slides: %s", order)

                # ── update_table ───────────────────────────────────────────
                elif op == "update_table":
                    slide_idx = rev.get("slide_index", 0)
                    row_idx   = rev.get("row", 0)
                    col_idx   = rev.get("col", 0)
                    value     = str(rev.get("value", ""))
                    if 0 <= slide_idx < len(prs.slides):
                        slide = prs.slides[slide_idx]
                        for shape in slide.shapes:
                            if shape.has_table:
                                tbl = shape.table
                                if row_idx < len(tbl.rows) and col_idx < len(tbl.rows[row_idx].cells):
                                    cell = tbl.rows[row_idx].cells[col_idx]
                                    if cell.text_frame.paragraphs:
                                        para = cell.text_frame.paragraphs[0]
                                        for run in para.runs:
                                            run.text = ""
                                        if para.runs:
                                            para.runs[0].text = value
                                        else:
                                            para.clear()
                                            run = para.add_run()
                                            run.text = value
                                    operations_applied += 1
                                break

                # ── append_slides ──────────────────────────────────────────
                # For append_slides, we generate the new slides using PptxSandbox
                # and then copy them into this presentation.
                elif op == "append_slides":
                    new_slides = rev.get("slides", [])
                    if new_slides:
                        try:
                            appended = self._append_slides(prs, new_slides)
                            operations_applied += appended
                        except Exception as ae:
                            logger.warning("append_slides error: %s", ae)

                # ── update_slide ───────────────────────────────────────────
                # Generates a new single-slide pptx and replaces the target slide
                elif op == "update_slide":
                    slide_idx  = rev.get("slide_index", 0)
                    slide_spec = rev.get("slide")
                    if slide_spec and 0 <= slide_idx < len(prs.slides):
                        try:
                            self._replace_slide_content(prs, slide_idx, slide_spec)
                            operations_applied += 1
                        except Exception as ue:
                            logger.warning("update_slide error at index %d: %s", slide_idx, ue)

            # Serialize
            buf = io.BytesIO()
            prs.save(buf)
            data = buf.getvalue()
            elapsed = round((time.time() - start) * 1000, 2)

            fname = output_filename or "revised_presentation.pptx"
            logger.info("PptxReviser ✓  ops=%d  replacements=%d  %.0f ms",
                        operations_applied, replacements_made, elapsed)

            return ReviseResult(
                success=True,
                data=data,
                filename=fname,
                operations_applied=operations_applied,
                replacements_made=replacements_made,
                exec_time_ms=elapsed,
            )

        except Exception as exc:
            logger.error("PptxReviser error: %s", exc, exc_info=True)
            return ReviseResult(
                success=False,
                error=str(exc),
                exec_time_ms=round((time.time() - start) * 1000, 2),
            )

    # =========================================================================
    # Helpers
    # =========================================================================

    @staticmethod
    def _find_replace_all(prs, find: str, replace: str) -> int:
        """Replace all occurrences of `find` with `replace` across all slides."""
        count = 0
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        for run in para.runs:
                            if find in run.text:
                                count += run.text.count(find)
                                run.text = run.text.replace(find, replace)
                if shape.has_table:
                    for row in shape.table.rows:
                        for cell in row.cells:
                            for para in cell.text_frame.paragraphs:
                                for run in para.runs:
                                    if find in run.text:
                                        count += run.text.count(find)
                                        run.text = run.text.replace(find, replace)
        return count

    @staticmethod
    def _replace_slide_content(prs, slide_idx: int, slide_spec: Dict[str, Any]) -> None:
        """
        Clear text content of a slide and write new text blocks from slide_spec.
        This does not change the slide layout — it only updates text.
        """
        slide = prs.slides[slide_idx]

        # Clear all text frames
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        run.text = ""

        # Find the first text frame and write the new title + content
        title_set = False
        content_set = False

        title_text = slide_spec.get("title", "")
        bullets = slide_spec.get("bullets", [])
        text = slide_spec.get("text", "")
        metrics = slide_spec.get("metrics", [])

        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            tf = shape.text_frame
            if shape.name.lower().startswith("title") and not title_set and title_text:
                if tf.paragraphs:
                    if tf.paragraphs[0].runs:
                        tf.paragraphs[0].runs[0].text = title_text.upper()
                    title_set = True
            elif not content_set:
                if bullets:
                    if tf.paragraphs:
                        if tf.paragraphs[0].runs:
                            tf.paragraphs[0].runs[0].text = "\n".join(f"• {b}" for b in bullets)
                    content_set = True
                elif text:
                    if tf.paragraphs:
                        if tf.paragraphs[0].runs:
                            tf.paragraphs[0].runs[0].text = text
                    content_set = True
                elif metrics:
                    metric_text = "  |  ".join(f"{m.get('value','')} {m.get('label','')}" for m in metrics)
                    if tf.paragraphs:
                        if tf.paragraphs[0].runs:
                            tf.paragraphs[0].runs[0].text = metric_text
                    content_set = True

    def _append_slides(self, prs, new_slides: List[Dict[str, Any]]) -> int:
        """
        Generate new slides using PptxSandbox and copy them into prs.
        Returns the number of slides appended.
        """
        try:
            from core.sandbox.pptx_sandbox import PptxSandbox

            temp_spec = {"title": "APPEND", "slides": new_slides}
            sandbox = PptxSandbox()
            result = sandbox.generate(temp_spec, timeout=60)
            if not result.success or not result.data:
                logger.warning("_append_slides: PptxSandbox generation failed: %s", result.error)
                return 0

            from pptx import Presentation as Prs
            new_prs = Prs(io.BytesIO(result.data))

            # Copy slides from new_prs into prs
            appended = 0
            for slide in new_prs.slides:
                # Add blank slide of same layout
                slide_layout = prs.slide_layouts[6]  # blank
                new_slide = prs.slides.add_slide(slide_layout)

                # Copy all shapes from the generated slide
                for shape in slide.shapes:
                    el = shape.element
                    new_slide.shapes._spTree.insert(2, copy.deepcopy(el))

                appended += 1

            return appended

        except Exception as exc:
            logger.warning("_append_slides error: %s", exc)
            return 0
