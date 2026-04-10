"""
PPTX Analyzer
=============
Reads any uploaded .pptx file and returns a structured profile.
Uses python-pptx (pure Python, no subprocess, instant execution).

The analyzer extracts:
- Slide count, titles, content types
- All text across the entire deck (flat + per-slide)
- Image locations
- Table data
- Potomac brand compliance score

Usage
-----
    from core.sandbox.pptx_analyzer import PptxAnalyzer

    analyzer = PptxAnalyzer()
    result   = analyzer.analyze(pptx_bytes, filename="deck.pptx")
    if result.success:
        print(result.profile)
"""

from __future__ import annotations

import io
import logging
import time
from dataclasses import dataclass, field
from typing import Any, Dict, List, Optional

logger = logging.getLogger(__name__)

# Potomac brand colors (no #)
_BRAND_COLORS = {"FEC00F", "212121", "FFFFFF", "FEF7D8", "FEE896", "999999", "DDDDDD", "F5F5F5", "EB2F5C", "00DED1", "22C55E"}


@dataclass
class SlideProfile:
    index: int
    title: str
    text_blocks: List[str]
    has_images: bool
    has_tables: bool
    table_data: List[List[List[str]]]  # list of tables, each table = list of rows, each row = list of cells
    shape_count: int


@dataclass
class AnalyzePptxResult:
    success: bool
    error: Optional[str] = None
    exec_time_ms: float = 0.0
    profile: Optional[Dict[str, Any]] = None


class PptxAnalyzer:
    """Read and profile any uploaded PowerPoint file."""

    def analyze(self, pptx_bytes: bytes, filename: Optional[str] = None) -> AnalyzePptxResult:
        start = time.time()
        try:
            from pptx import Presentation
            from pptx.util import Pt

            prs = Presentation(io.BytesIO(pptx_bytes))

            slides_profile = []
            all_text_flat  = []
            brand_violations = []

            for slide_idx, slide in enumerate(prs.slides):
                title_text   = ""
                text_blocks  = []
                has_images   = False
                has_tables   = False
                table_data   = []
                shape_count  = len(slide.shapes)

                for shape in slide.shapes:
                    stype = shape.shape_type

                    # Title
                    if shape.has_text_frame:
                        for para in shape.text_frame.paragraphs:
                            para_text = para.text.strip()
                            if not para_text:
                                continue
                            if shape.name.lower().startswith("title") or (slide_idx == 0 and not title_text):
                                if not title_text:
                                    title_text = para_text[:120]
                            text_blocks.append(para_text[:300])
                            all_text_flat.append(para_text[:300])

                            # Brand compliance: check font colors
                            for run in para.runs:
                                if run.font.color and run.font.color.type:
                                    try:
                                        rgb = run.font.color.rgb
                                        hex_color = str(rgb).upper()
                                        if hex_color not in _BRAND_COLORS:
                                            brand_violations.append({
                                                "slide": slide_idx + 1,
                                                "issue": f"Off-brand font color #{hex_color}",
                                                "shape": shape.name,
                                            })
                                    except Exception:
                                        pass

                    # Images
                    if stype == 13:  # MSO_SHAPE_TYPE.PICTURE
                        has_images = True

                    # Tables
                    if shape.has_table:
                        has_tables = True
                        tbl = shape.table
                        tbl_rows = []
                        for row in tbl.rows:
                            tbl_rows.append([cell.text.strip() for cell in row.cells])
                        table_data.append(tbl_rows)

                    # Background fill color check
                    try:
                        if shape.fill.type == 1:  # SOLID
                            rgb = str(shape.fill.fore_color.rgb).upper()
                            if rgb not in _BRAND_COLORS:
                                brand_violations.append({
                                    "slide": slide_idx + 1,
                                    "issue": f"Off-brand fill color #{rgb}",
                                    "shape": shape.name,
                                })
                    except Exception:
                        pass

                # Check for missing title
                if not title_text and text_blocks:
                    title_text = text_blocks[0][:80]

                slides_profile.append(SlideProfile(
                    index=slide_idx + 1,
                    title=title_text,
                    text_blocks=text_blocks[:20],
                    has_images=has_images,
                    has_tables=has_tables,
                    table_data=table_data,
                    shape_count=shape_count,
                ))

            # Brand compliance score
            max_violations = len(prs.slides) * 3
            violation_count = min(len(brand_violations), max_violations)
            compliance_score = max(0, round(100 - (violation_count / max(max_violations, 1)) * 100))

            profile = {
                "filename":     filename or "presentation.pptx",
                "slide_count":  len(prs.slides),
                "all_text":     " | ".join(all_text_flat[:200]),
                "slides":       [
                    {
                        "index":       sp.index,
                        "title":       sp.title,
                        "text_blocks": sp.text_blocks,
                        "has_images":  sp.has_images,
                        "has_tables":  sp.has_tables,
                        "table_data":  sp.table_data,
                        "shape_count": sp.shape_count,
                    }
                    for sp in slides_profile
                ],
                "brand_compliance": {
                    "score":      compliance_score,
                    "violations": brand_violations[:50],
                },
            }

            elapsed = round((time.time() - start) * 1000, 2)
            logger.info("PptxAnalyzer ✓  %d slides, compliance=%d%%, %.0f ms",
                        len(prs.slides), compliance_score, elapsed)

            return AnalyzePptxResult(success=True, profile=profile, exec_time_ms=elapsed)

        except Exception as exc:
            logger.error("PptxAnalyzer error: %s", exc, exc_info=True)
            return AnalyzePptxResult(
                success=False,
                error=str(exc),
                exec_time_ms=round((time.time() - start) * 1000, 2),
            )
