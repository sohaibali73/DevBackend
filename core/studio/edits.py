"""
Visual-editor JSON-ops applier for Studio artifacts.

The frontend sends a list of edit operations; we re-open the source artifact's
.pptx/.docx, apply ops via python-pptx / python-docx, then register a new
artifact version (v(n+1)).

Supported ops (frontend ↔ backend contract):

    PPTX:
      {"type":"text",   "slide": int (1-based), "shape_index": int, "value": "new text"}
      {"type":"text_replace", "slide": int, "find": "...", "replace": "...", "all": false}
      {"type":"add_slide_note", "slide": int, "value": "..."}
      {"type":"reorder_slides", "order": [3,1,2,4]}
      {"type":"delete_slide", "slide": int}
      {"type":"duplicate_slide", "slide": int}

    DOCX:
      {"type":"text_replace", "find": "...", "replace": "...", "all": true}
      {"type":"replace_paragraph", "index": int, "value": "..."}
      {"type":"append_paragraph", "value": "...", "style": "Normal"}
      {"type":"append_heading", "value": "...", "level": 1}

Unknown op types are skipped (with a warning) so future frontend additions
don't crash old backends.
"""

from __future__ import annotations

import io
import logging
from typing import Any, Dict, List, Optional, Tuple

logger = logging.getLogger(__name__)


# =============================================================================
# PPTX
# =============================================================================

def _apply_pptx_ops(prs_bytes: bytes, ops: List[Dict[str, Any]]) -> bytes:
    from pptx import Presentation
    from copy import deepcopy

    prs = Presentation(io.BytesIO(prs_bytes))

    def _slide_at(one_based: int):
        idx = one_based - 1
        if idx < 0 or idx >= len(prs.slides):
            raise IndexError(f"slide {one_based} out of range")
        return prs.slides[idx]

    for op in ops:
        try:
            t = op.get("type", "")
            if t == "text":
                slide = _slide_at(int(op["slide"]))
                shapes = list(slide.shapes)
                shp = shapes[int(op["shape_index"])]
                if shp.has_text_frame:
                    tf = shp.text_frame
                    # Replace text but preserve first run's formatting
                    new_text = str(op.get("value", ""))
                    if tf.paragraphs:
                        p = tf.paragraphs[0]
                        if p.runs:
                            p.runs[0].text = new_text
                            for r in p.runs[1:]:
                                r.text = ""
                        else:
                            p.text = new_text
                        # Clear remaining paragraphs
                        for p2 in tf.paragraphs[1:]:
                            for r in p2.runs:
                                r.text = ""
                    else:
                        tf.text = new_text

            elif t == "text_replace":
                find = str(op.get("find", ""))
                repl = str(op.get("replace", ""))
                target_slide = op.get("slide")
                replace_all = bool(op.get("all", target_slide is None))
                slides_iter = (
                    [_slide_at(int(target_slide))]
                    if target_slide is not None
                    else list(prs.slides)
                )
                for s in slides_iter:
                    for shp in s.shapes:
                        if not shp.has_text_frame:
                            continue
                        for para in shp.text_frame.paragraphs:
                            for run in para.runs:
                                if find and find in (run.text or ""):
                                    run.text = run.text.replace(find, repl)
                                    if not replace_all:
                                        break

            elif t == "add_slide_note":
                slide = _slide_at(int(op["slide"]))
                notes = slide.notes_slide.notes_text_frame
                notes.text = str(op.get("value", ""))

            elif t == "reorder_slides":
                order = [int(x) - 1 for x in op.get("order", [])]
                if sorted(order) != list(range(len(prs.slides))):
                    logger.warning("reorder_slides: order must be a permutation; skipping")
                    continue
                xml_slides = prs.slides._sldIdLst  # noqa: SLF001
                slides_xml = list(xml_slides)
                xml_slides.clear()
                for i in order:
                    xml_slides.append(slides_xml[i])

            elif t == "delete_slide":
                idx = int(op["slide"]) - 1
                xml_slides = prs.slides._sldIdLst  # noqa: SLF001
                slides_xml = list(xml_slides)
                if 0 <= idx < len(slides_xml):
                    xml_slides.remove(slides_xml[idx])

            elif t == "duplicate_slide":
                # Cheap duplicate via XML clone; layout/relationships are preserved
                idx = int(op["slide"]) - 1
                src = prs.slides[idx]
                blank_layout = src.slide_layout
                new_slide = prs.slides.add_slide(blank_layout)
                # Copy shapes (best-effort; complex shapes like charts may not duplicate cleanly)
                for shape in src.shapes:
                    el = shape.element
                    new_el = deepcopy(el)
                    new_slide.shapes._spTree.insert_element_before(new_el, "p:extLst")  # noqa: SLF001

            else:
                logger.info("Unknown PPTX op skipped: %s", t)

        except Exception as e:
            logger.warning("PPTX op %s failed: %s", op, e)

    out = io.BytesIO()
    prs.save(out)
    return out.getvalue()


# =============================================================================
# DOCX
# =============================================================================

def _apply_docx_ops(doc_bytes: bytes, ops: List[Dict[str, Any]]) -> bytes:
    from docx import Document

    doc = Document(io.BytesIO(doc_bytes))

    for op in ops:
        try:
            t = op.get("type", "")
            if t == "text_replace":
                find = str(op.get("find", ""))
                repl = str(op.get("replace", ""))
                replace_all = bool(op.get("all", True))
                done = False
                for para in doc.paragraphs:
                    if find and find in para.text:
                        # Re-write the paragraph; preserves first run's formatting
                        if para.runs:
                            new_text = para.text.replace(find, repl)
                            para.runs[0].text = new_text
                            for r in para.runs[1:]:
                                r.text = ""
                        if not replace_all:
                            done = True
                            break
                    if done:
                        break

            elif t == "replace_paragraph":
                idx = int(op["index"])
                if 0 <= idx < len(doc.paragraphs):
                    p = doc.paragraphs[idx]
                    val = str(op.get("value", ""))
                    if p.runs:
                        p.runs[0].text = val
                        for r in p.runs[1:]:
                            r.text = ""
                    else:
                        p.text = val

            elif t == "append_paragraph":
                style = op.get("style") or "Normal"
                doc.add_paragraph(str(op.get("value", "")), style=style)

            elif t == "append_heading":
                doc.add_heading(str(op.get("value", "")), level=int(op.get("level", 1)))

            else:
                logger.info("Unknown DOCX op skipped: %s", t)

        except Exception as e:
            logger.warning("DOCX op %s failed: %s", op, e)

    out = io.BytesIO()
    doc.save(out)
    return out.getvalue()


# =============================================================================
# Public entry — apply ops to an artifact and persist new version
# =============================================================================

def apply_ops_to_artifact(
    *,
    user_id: str,
    artifact_id: str,
    ops: List[Dict[str, Any]],
    save_edit_state: bool = True,
) -> Dict[str, Any]:
    """
    Load the artifact bytes, apply ops, register a new artifact version.
    Returns the new artifact row (dict).
    """
    from core.studio.projects import (
        artifact_bytes,
        get_artifact,
        register_artifact_from_bytes,
    )

    src = get_artifact(artifact_id, user_id)
    if not src:
        raise ValueError("artifact not found")

    data, filename, _mime = artifact_bytes(artifact_id, user_id)
    if not data:
        raise FileNotFoundError("artifact bytes missing on volume")

    kind = src["kind"]
    if kind == "pptx":
        new_bytes = _apply_pptx_ops(data, ops or [])
    elif kind == "docx":
        new_bytes = _apply_docx_ops(data, ops or [])
    else:
        raise ValueError(f"Unsupported artifact kind for editing: {kind}")

    return register_artifact_from_bytes(
        user_id=user_id,
        project_id=src["project_id"],
        kind=kind,
        data=new_bytes,
        filename=filename or f"v{src['version']+1}.{kind}",
        conversation_id=src.get("conversation_id"),
        message_id=None,
        source_file_id=src.get("source_file_id"),
        edit_state=({"ops": ops, "from_version": src["version"]} if save_edit_state else None),
        meta={
            "edited_from_artifact": src["id"],
            "edit_op_count":        len(ops or []),
        },
    )
