#!/usr/bin/env python3
"""
local_parser.py — Client-side document text extractor for the KB Uploader GUI.

Mirrors the server's DocumentParser so text is extracted locally before upload.
The server then skips parsing entirely (use /kb-admin/upload-preparsed).

Supported formats: PDF, DOCX, DOC, XLSX, XLS, PPTX, CSV, TXT, MD, HTML, JSON, XML, RTF
Uses the same libraries already in requirements.txt:
    pypdf, PyMuPDF, python-docx, openpyxl, python-pptx
"""

import csv
import hashlib
import io
import json
import logging
import re
import os
from pathlib import Path
from typing import Tuple, NamedTuple

logger = logging.getLogger(__name__)


# ─────────────────────────────────────────────────────────────────────────────

class ParseResult(NamedTuple):
    text: str          # extracted + cleaned text
    error: str         # "" on success, error message on failure
    truncated: bool    # True if text was capped


MAX_CHARS = 200_000    # ~40 k tokens — enough for RAG, avoids memory blowup

# ─────────────────────────────────────────────────────────────────────────────
# PUBLIC API
# ─────────────────────────────────────────────────────────────────────────────

def parse_file(path: Path) -> ParseResult:
    """
    Extract text from `path`.  Never raises.
    Returns ParseResult(text, error, truncated).
    """
    ext = path.suffix.lower()
    try:
        if ext == ".pdf":
            raw, err = _parse_pdf(path)
        elif ext == ".docx":
            raw, err = _parse_docx(path)
        elif ext == ".doc":
            raw, err = _parse_doc(path)
        elif ext in (".xlsx", ".xls"):
            raw, err = _parse_excel(path)
        elif ext == ".pptx":
            raw, err = _parse_pptx(path)
        elif ext == ".csv":
            raw, err = _parse_csv(path)
        elif ext in (".html", ".htm"):
            raw, err = _parse_html(path)
        elif ext == ".json":
            raw, err = _parse_json(path)
        elif ext == ".xml":
            raw, err = _parse_xml(path)
        elif ext == ".rtf":
            raw, err = _parse_rtf(path)
        else:
            # Plain text fallback (txt, md, log, afl, py, …)
            raw = path.read_text(encoding="utf-8", errors="ignore")
            err = ""
    except Exception as exc:
        return ParseResult("", str(exc), False)

    if err:
        return ParseResult("", err, False)

    cleaned = _clean(raw)
    truncated = len(cleaned) > MAX_CHARS
    return ParseResult(cleaned[:MAX_CHARS], "", truncated)


def file_hash(path: Path) -> str:
    """SHA-256 of the raw file bytes (matches server-side dedup)."""
    h = hashlib.sha256()
    with open(path, "rb") as f:
        for chunk in iter(lambda: f.read(65536), b""):
            h.update(chunk)
    return h.hexdigest()


def available_parsers() -> dict:
    """Return a dict of format → available library name (for diagnostics)."""
    avail = {}
    try:
        import fitz; avail["PDF (fast)"] = f"PyMuPDF {fitz.__version__}"
    except ImportError:
        pass
    try:
        import pypdf; avail["PDF (fallback)"] = "pypdf"
    except ImportError:
        pass
    try:
        import docx; avail["DOCX"] = "python-docx"
    except ImportError:
        pass
    try:
        import openpyxl; avail["XLSX"] = f"openpyxl {openpyxl.__version__}"
    except ImportError:
        pass
    try:
        import pptx; avail["PPTX"] = "python-pptx"
    except ImportError:
        pass
    return avail


# ─────────────────────────────────────────────────────────────────────────────
# CLEANERS
# ─────────────────────────────────────────────────────────────────────────────

_CTRL  = re.compile(r'[\x00-\x08\x0B\x0C\x0E-\x1F\x7F]')
_PAGES = re.compile(r'Page\s+\d+\s+of\s+\d+', re.IGNORECASE)
_NL3   = re.compile(r'\n{3,}')
_SP2   = re.compile(r' {2,}')


def _clean(text: str) -> str:
    text = _CTRL.sub('', text)
    text = _PAGES.sub('', text)
    text = _NL3.sub('\n\n', text)
    text = _SP2.sub(' ', text)
    return text.strip()


# ─────────────────────────────────────────────────────────────────────────────
# FORMAT PARSERS
# ─────────────────────────────────────────────────────────────────────────────

def _parse_pdf(path: Path) -> Tuple[str, str]:
    # ── PyMuPDF (fastest) ──────────────────────────────────────────────────
    try:
        import fitz
        doc = fitz.open(str(path))
        text = "".join(page.get_text("text") for page in doc)
        doc.close()
        if text.strip():
            return text, ""
    except ImportError:
        pass
    except Exception as exc:
        logger.debug(f"PyMuPDF failed ({path.name}): {exc}")

    # ── pypdf ──────────────────────────────────────────────────────────────
    try:
        from pypdf import PdfReader
        reader = PdfReader(str(path))
        text = "\n".join(page.extract_text() or "" for page in reader.pages)
        if text.strip():
            return text, ""
    except ImportError:
        pass
    except Exception as exc:
        logger.debug(f"pypdf failed ({path.name}): {exc}")

    # ── pdfplumber ────────────────────────────────────────────────────────
    try:
        import pdfplumber
        with pdfplumber.open(str(path)) as pdf:
            text = "\n".join(p.extract_text() or "" for p in pdf.pages)
        return text, ""
    except ImportError:
        return "", "PDF parsing unavailable — install:  pip install pypdf PyMuPDF"
    except Exception as exc:
        return "", f"PDF error: {exc}"


def _parse_docx(path: Path) -> Tuple[str, str]:
    try:
        from docx import Document
        doc = Document(str(path))
        parts = []

        # Paragraphs
        for p in doc.paragraphs:
            if p.text.strip():
                parts.append(p.text.strip())

        # Tables
        for table in doc.tables:
            for row in table.rows:
                cells = [c.text.strip() for c in row.cells if c.text.strip()]
                if cells:
                    parts.append(" | ".join(cells))

        return "\n".join(parts), ""
    except ImportError:
        return "", "DOCX parsing unavailable — install:  pip install python-docx"
    except Exception as exc:
        return "", f"DOCX error: {exc}"


def _parse_doc(path: Path) -> Tuple[str, str]:
    """Legacy .doc — crude ASCII text extraction (no antiword required)."""
    try:
        with open(path, "rb") as f:
            raw = f.read()
        text = re.sub(rb'[^\x20-\x7E\n\r\t]', b' ', raw).decode("ascii", errors="ignore")
        text = re.sub(r' {4,}', ' ', text)
        return text, ""
    except Exception as exc:
        return "", f"DOC error: {exc}"


def _parse_excel(path: Path) -> Tuple[str, str]:
    # ── openpyxl (.xlsx) ──────────────────────────────────────────────────
    try:
        import openpyxl
        wb = openpyxl.load_workbook(str(path), read_only=True, data_only=True)
        parts = []
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            rows = []
            for i, row in enumerate(ws.iter_rows(max_row=1000, values_only=True)):
                row_text = " | ".join(
                    str(v) for v in row if v is not None and str(v).strip()
                )
                if row_text:
                    rows.append(row_text)
                if i >= 999:
                    rows.append("… (truncated at 1000 rows)")
                    break
            if rows:
                parts.append(f"=== Sheet: {sheet_name} ===\n" + "\n".join(rows))
        wb.close()
        return "\n\n".join(parts), ""
    except ImportError:
        pass
    except Exception as exc:
        logger.debug(f"openpyxl failed ({path.name}): {exc}")

    # ── xlrd (.xls) ───────────────────────────────────────────────────────
    try:
        import xlrd
        wb = xlrd.open_workbook(str(path))
        parts = []
        for sheet in wb.sheets():
            rows = []
            for i in range(min(sheet.nrows, 1000)):
                row_text = " | ".join(str(v) for v in sheet.row_values(i) if v)
                if row_text:
                    rows.append(row_text)
            if rows:
                parts.append(f"=== Sheet: {sheet.name} ===\n" + "\n".join(rows))
        return "\n\n".join(parts), ""
    except ImportError:
        return "", "Excel parsing unavailable — install:  pip install openpyxl"
    except Exception as exc:
        return "", f"Excel error: {exc}"


def _parse_pptx(path: Path) -> Tuple[str, str]:
    try:
        from pptx import Presentation
        prs = Presentation(str(path))
        parts = []
        for i, slide in enumerate(prs.slides, 1):
            texts = [
                shape.text.strip()
                for shape in slide.shapes
                if hasattr(shape, "text") and shape.text.strip()
            ]
            if texts:
                parts.append(f"=== Slide {i} ===\n" + "\n".join(texts))
        return "\n\n".join(parts) or "[Empty presentation]", ""
    except ImportError:
        return "", "PPTX parsing unavailable — install:  pip install python-pptx"
    except Exception as exc:
        return "", f"PPTX error: {exc}"


def _parse_csv(path: Path) -> Tuple[str, str]:
    try:
        rows = []
        with open(path, newline="", encoding="utf-8", errors="ignore") as f:
            reader = csv.reader(f)
            for i, row in enumerate(reader):
                if i >= 1000:
                    rows.append("… (truncated at 1000 rows)")
                    break
                rows.append(" | ".join(str(v) for v in row))
        return "\n".join(rows), ""
    except Exception as exc:
        return "", f"CSV error: {exc}"


def _parse_html(path: Path) -> Tuple[str, str]:
    try:
        from html.parser import HTMLParser

        class _Extractor(HTMLParser):
            def __init__(self):
                super().__init__()
                self.parts: list = []
                self._skip = False

            def handle_starttag(self, tag, attrs):
                if tag in ("script", "style", "head"):
                    self._skip = True
                elif tag in ("p", "div", "br", "li", "h1", "h2", "h3",
                             "h4", "h5", "h6", "tr", "td", "th"):
                    self.parts.append("\n")

            def handle_endtag(self, tag):
                if tag in ("script", "style", "head"):
                    self._skip = False

            def handle_data(self, data):
                if not self._skip and data.strip():
                    self.parts.append(data.strip())

        extractor = _Extractor()
        extractor.feed(path.read_text(encoding="utf-8", errors="ignore"))
        return " ".join(extractor.parts), ""
    except Exception as exc:
        return "", f"HTML error: {exc}"


def _parse_json(path: Path) -> Tuple[str, str]:
    try:
        with open(path, encoding="utf-8") as f:
            data = json.load(f)
        return json.dumps(data, indent=2, ensure_ascii=False)[:80000], ""
    except Exception as exc:
        # Return raw content if not valid JSON
        try:
            return path.read_text(encoding="utf-8", errors="ignore"), ""
        except Exception:
            return "", f"JSON error: {exc}"


def _parse_xml(path: Path) -> Tuple[str, str]:
    try:
        raw = path.read_text(encoding="utf-8", errors="ignore")
        text = re.sub(r'<[^>]+>', ' ', raw)
        text = re.sub(r'&\w+;', ' ', text)
        return text, ""
    except Exception as exc:
        return "", f"XML error: {exc}"


def _parse_rtf(path: Path) -> Tuple[str, str]:
    """Crude RTF → plain text (strips control words)."""
    try:
        raw = path.read_text(encoding="utf-8", errors="ignore")
        text = re.sub(r'\\[a-z]+\d* ?', '', raw)
        text = re.sub(r'[{}]', '', text)
        return text, ""
    except Exception as exc:
        return "", f"RTF error: {exc}"
