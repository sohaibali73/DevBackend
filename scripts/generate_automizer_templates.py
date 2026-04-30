"""
generate_automizer_templates.py
================================
Generates Potomac-branded .pptx template files for the generate_pptx_template
tool (pptx-automizer integration).

Each slide has NAMED shapes (visible in PowerPoint Selection Pane via ALT+F10)
so pptx-automizer can target them with:
  - set_text(shape, text)
  - replace_tagged(shape, tags)
  - set_chart_data(shape, series, categories)
  - set_table(shape, body)
  - swap_image(shape, image_file)

Output: ClaudeSkills/potomac-pptx/automizer-templates/
  - potomac-content-slides.pptx   (5 slides: title, content, two-col, metrics, divider)
  - potomac-chart-slides.pptx     (3 slides: bar chart, line chart, column/attribution)
  - potomac-table-slides.pptx     (3 slides: data table, holdings table, scorecard table)
  - potomac-fund-fact-sheet.pptx  (1 slide: complete fund fact sheet with {{tags}})

Usage:
    python scripts/generate_automizer_templates.py
"""

from __future__ import annotations

import sys
from pathlib import Path

# ── Project root on path ──────────────────────────────────────────────────────
_PROJECT_ROOT = Path(__file__).parent.parent
sys.path.insert(0, str(_PROJECT_ROOT))

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.chart.data import CategoryChartData
    from pptx.enum.chart import XL_CHART_TYPE
except ImportError:
    print("ERROR: python-pptx not installed. Run: pip install python-pptx")
    sys.exit(1)

# ── Potomac brand palette ─────────────────────────────────────────────────────
YELLOW    = RGBColor(0xFE, 0xC0, 0x0F)
DARK_GRAY = RGBColor(0x21, 0x21, 0x21)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
GRAY_60   = RGBColor(0x99, 0x99, 0x99)
GRAY_20   = RGBColor(0xDD, 0xDD, 0xDD)
YELLOW_20 = RGBColor(0xFE, 0xF7, 0xD8)
GREEN     = RGBColor(0x22, 0xC5, 0x5E)
RED       = RGBColor(0xEB, 0x2F, 0x5C)

# ── Potomac brand fonts ───────────────────────────────────────────────────────
FONT_H = "Rajdhani"    # Headline — always ALL CAPS
FONT_B = "Quicksand"   # Body / captions

# ── Slide dimensions: 10" × 7.5" (LAYOUT_WIDE) ───────────────────────────────
SLIDE_W = Inches(10)
SLIDE_H = Inches(7.5)

# ── MSO auto-shape type integer for rectangle ─────────────────────────────────
RECT = 1  # MSO_AUTO_SHAPE_TYPE.RECTANGLE

# ── Output directory ──────────────────────────────────────────────────────────
OUTPUT_DIR = _PROJECT_ROOT / "ClaudeSkills" / "potomac-pptx" / "automizer-templates"
OUTPUT_DIR.mkdir(parents=True, exist_ok=True)


# =============================================================================
# Low-level helpers
# =============================================================================

def new_prs() -> Presentation:
    """Return a fresh LAYOUT_WIDE (10x7.5") presentation."""
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def blank_slide(prs: Presentation):
    """Add and return a completely blank slide."""
    return prs.slides.add_slide(prs.slide_layouts[6])


def rect(slide, name: str, x, y, w, h,
         fill: RGBColor = None, line: RGBColor = None):
    """Add a named rectangle. Pass fill=None for no-fill, line=None for no border."""
    shp = slide.shapes.add_shape(RECT, x, y, w, h)
    shp.name = name
    if fill:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
    else:
        shp.fill.background()
    if line:
        shp.line.color.rgb = line
    else:
        shp.line.fill.background()
    return shp


def txbox(slide, name: str, x, y, w, h,
          text: str = "",
          font: str = FONT_B, size: int = 16,
          color: RGBColor = DARK_GRAY,
          bold: bool = False,
          align: PP_ALIGN = PP_ALIGN.LEFT):
    """Add a named text box with Potomac-styled text."""
    box = slide.shapes.add_textbox(x, y, w, h)
    box.name = name
    tf = box.text_frame
    tf.word_wrap = True
    para = tf.paragraphs[0]
    para.alignment = align
    run = para.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def add_header(slide, logo_text: str = "POTOMAC"):
    """Standard Potomac header: yellow accent bar + logo text."""
    rect(slide, "AccentBar",
         Inches(0), Inches(0), Inches(0.15), Inches(7.5),
         fill=YELLOW)
    txbox(slide, "LogoArea",
          Inches(8.55), Inches(0.15), Inches(1.25), Inches(0.5),
          text=logo_text,
          font=FONT_H, size=12, color=DARK_GRAY, bold=True,
          align=PP_ALIGN.CENTER)


def add_title_underline(slide):
    """Short yellow underline below the slide title."""
    rect(slide, "TitleUnderline",
         Inches(0.5), Inches(1.3), Inches(2.5), Inches(0.06),
         fill=YELLOW)


def add_slide_title(slide, placeholder: str = "SLIDE TITLE",
                    font_size: int = 26):
    """Standard title text box for content slides."""
    txbox(slide, "TitleText",
          Inches(0.5), Inches(0.35), Inches(7.9), Inches(0.9),
          text=placeholder, font=FONT_H, size=font_size,
          color=DARK_GRAY, bold=True)
    add_title_underline(slide)


def chart_slide_setup(prs, title_placeholder: str):
    """Return a new slide with header + title already applied."""
    slide = blank_slide(prs)
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE
    add_header(slide)
    add_slide_title(slide, placeholder=title_placeholder)
    return slide


def add_caption(slide, placeholder: str = "Source: Potomac | As of {{date}}"):
    txbox(slide, "CaptionText",
          Inches(0.5), Inches(7.1), Inches(9), Inches(0.3),
          text=placeholder, font=FONT_B, size=10,
          color=GRAY_60, align=PP_ALIGN.CENTER)


# =============================================================================
# Table helper — creates a properly styled Potomac table
# =============================================================================

def add_table(slide, name: str, headers: list, rows: list,
              col_widths: list, x=Inches(0.5), y=Inches(1.5)):
    """
    Add a named, Potomac-branded table.
    - Row 0  : header (yellow background, FONT_H)
    - Rows 1+: data (zebra fill, FONT_B)
    """
    n_rows = len(rows) + 1   # +1 for header
    n_cols = len(headers)
    total_w = sum(col_widths)
    row_h   = Inches(0.44)

    tbl_shape = slide.shapes.add_table(n_rows, n_cols, x, y,
                                        total_w, row_h * n_rows)
    tbl_shape.name = name
    tbl = tbl_shape.table

    # Set column widths
    for i, cw in enumerate(col_widths):
        tbl.columns[i].width = cw

    # Header row
    for j, hdr in enumerate(headers):
        cell = tbl.cell(0, j)
        cell.fill.solid()
        cell.fill.fore_color.rgb = YELLOW
        para = cell.text_frame.paragraphs[0]
        para.alignment = PP_ALIGN.CENTER
        run  = para.add_run()
        run.text = hdr.upper()
        run.font.name = FONT_H
        run.font.size = Pt(11)
        run.font.bold = True
        run.font.color.rgb = DARK_GRAY

    # Data rows
    for i, row_data in enumerate(rows):
        bg = GRAY_20 if i % 2 == 1 else WHITE
        for j, val in enumerate(row_data):
            cell = tbl.cell(i + 1, j)
            cell.fill.solid()
            cell.fill.fore_color.rgb = bg
            para = cell.text_frame.paragraphs[0]
            para.alignment = PP_ALIGN.CENTER if j > 0 else PP_ALIGN.LEFT
            run  = para.add_run()
            run.text = str(val)
            run.font.name = FONT_B
            run.font.size = Pt(11)
            run.font.color.rgb = DARK_GRAY

    return tbl_shape


# =============================================================================
# Template 1: potomac-content-slides.pptx
# =============================================================================

def build_content_slides():
    """
    5 slides:
      1. Title slide   → TitleText, SubtitleText, DateText, TaglineText
      2. Content slide → TitleText, BodyText
      3. Two-column    → TitleText, LeftHeader, RightHeader, LeftContent, RightContent
      4. Metrics       → TitleText, Metric1-3 Value+Label, ContextText
      5. Section divider → TitleText, DescriptionText
    """
    prs = new_prs()

    # ── Slide 1: Title ────────────────────────────────────────────────────────
    slide = blank_slide(prs)
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE

    txbox(slide, "LogoArea",
          Inches(4.1), Inches(0.45), Inches(1.8), Inches(0.9),
          text="POTOMAC", font=FONT_H, size=20, color=DARK_GRAY,
          bold=True, align=PP_ALIGN.CENTER)

    txbox(slide, "TitleText",
          Inches(0.5), Inches(2.1), Inches(9), Inches(1.5),
          text="PRESENTATION TITLE",
          font=FONT_H, size=44, color=DARK_GRAY, bold=True,
          align=PP_ALIGN.CENTER)

    txbox(slide, "SubtitleText",
          Inches(0.5), Inches(3.8), Inches(9), Inches(0.8),
          text="Subtitle — {{subtitle}}",
          font=FONT_B, size=20, color=GRAY_60,
          align=PP_ALIGN.CENTER)

    rect(slide, "AccentBar",
         Inches(0.5), Inches(5.2), Inches(9), Inches(0.08), fill=YELLOW)

    txbox(slide, "DateText",
          Inches(0.5), Inches(5.4), Inches(9), Inches(0.5),
          text="{{date}}",
          font=FONT_B, size=14, color=GRAY_60, align=PP_ALIGN.CENTER)

    txbox(slide, "TaglineText",
          Inches(0.5), Inches(6.0), Inches(9), Inches(0.5),
          text="",
          font=FONT_B, size=15, color=YELLOW, align=PP_ALIGN.CENTER)

    # ── Slide 2: Content ──────────────────────────────────────────────────────
    slide = blank_slide(prs)
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE
    add_header(slide)
    add_slide_title(slide, "SLIDE TITLE", font_size=30)

    txbox(slide, "BodyText",
          Inches(0.5), Inches(1.55), Inches(9.1), Inches(5.5),
          text="\u2022 Key point one\n\u2022 Key point two\n\u2022 Key point three",
          font=FONT_B, size=18, color=DARK_GRAY)

    # ── Slide 3: Two-Column ───────────────────────────────────────────────────
    slide = blank_slide(prs)
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE
    add_header(slide)
    add_slide_title(slide, "TWO COLUMN SLIDE", font_size=26)

    txbox(slide, "LeftHeader",
          Inches(0.5), Inches(1.55), Inches(4.3), Inches(0.42),
          text="LEFT COLUMN",
          font=FONT_B, size=14, color=YELLOW, bold=True)
    txbox(slide, "RightHeader",
          Inches(5.35), Inches(1.55), Inches(4.3), Inches(0.42),
          text="RIGHT COLUMN",
          font=FONT_B, size=14, color=YELLOW, bold=True)

    rect(slide, "ColumnDivider",
         Inches(4.88), Inches(1.55), Inches(0.04), Inches(5.5), fill=GRAY_20)

    txbox(slide, "LeftContent",
          Inches(0.5), Inches(2.1), Inches(4.3), Inches(4.9),
          text="Left column content goes here.",
          font=FONT_B, size=15, color=DARK_GRAY)
    txbox(slide, "RightContent",
          Inches(5.35), Inches(2.1), Inches(4.3), Inches(4.9),
          text="Right column content goes here.",
          font=FONT_B, size=15, color=DARK_GRAY)

    # ── Slide 4: Metrics ──────────────────────────────────────────────────────
    slide = blank_slide(prs)
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE
    add_header(slide)
    add_slide_title(slide, "KEY METRICS", font_size=28)

    metric_cfg = [
        ("Metric1Value", "Metric1Label", Inches(0.6)),
        ("Metric2Value", "Metric2Label", Inches(3.73)),
        ("Metric3Value", "Metric3Label", Inches(6.86)),
    ]
    for val_name, lbl_name, x_pos in metric_cfg:
        txbox(slide, val_name,
              x_pos, Inches(2.0), Inches(2.7), Inches(1.2),
              text="{{value}}",
              font=FONT_H, size=54, color=YELLOW, bold=True,
              align=PP_ALIGN.CENTER)
        txbox(slide, lbl_name,
              x_pos, Inches(3.2), Inches(2.7), Inches(0.6),
              text="Metric Label",
              font=FONT_B, size=14, color=GRAY_60,
              align=PP_ALIGN.CENTER)

    txbox(slide, "ContextText",
          Inches(0.5), Inches(6.5), Inches(9), Inches(0.75),
          text="Source: {{source}} | As of {{date}}",
          font=FONT_B, size=11, color=GRAY_60, align=PP_ALIGN.CENTER)

    # ── Slide 5: Section Divider ──────────────────────────────────────────────
    slide = blank_slide(prs)
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = YELLOW_20

    txbox(slide, "LogoArea",
          Inches(8.55), Inches(0.15), Inches(1.25), Inches(0.5),
          text="POTOMAC", font=FONT_H, size=12, color=DARK_GRAY,
          bold=True, align=PP_ALIGN.CENTER)

    rect(slide, "AccentBar",
         Inches(0), Inches(0), Inches(0.45), Inches(7.5), fill=YELLOW)

    txbox(slide, "TitleText",
          Inches(0.8), Inches(2.4), Inches(8.0), Inches(1.6),
          text="SECTION TITLE",
          font=FONT_H, size=42, color=DARK_GRAY, bold=True)

    txbox(slide, "DescriptionText",
          Inches(0.8), Inches(4.2), Inches(8.0), Inches(1.5),
          text="Section description — {{description}}",
          font=FONT_B, size=18, color=GRAY_60)

    out = OUTPUT_DIR / "potomac-content-slides.pptx"
    prs.save(str(out))
    print(f"  \u2713 {out.name}")


# =============================================================================
# Template 2: potomac-chart-slides.pptx
# =============================================================================

def _make_chart_slide(prs, title: str, chart_name: str,
                      chart_type, categories: list, series_data: list):
    """Helper that adds one chart slide."""
    slide = chart_slide_setup(prs, title)

    chart_data = CategoryChartData()
    chart_data.categories = categories
    for s in series_data:
        chart_data.add_series(s["name"], s["values"])

    gf = slide.shapes.add_chart(
        chart_type,
        Inches(0.5), Inches(1.5),
        Inches(9.0), Inches(5.5),
        chart_data,
    )
    gf.name = chart_name
    gf.chart.has_legend = len(series_data) > 1

    add_caption(slide)
    return slide


def build_chart_slides():
    """
    3 slides:
      1. Performance bar chart   → TitleText, PerformanceChart,  CaptionText
      2. Cumulative line chart   → TitleText, TrendChart,         CaptionText
      3. Attribution column chart → TitleText, AttributionChart,  CaptionText
    """
    prs = new_prs()

    _make_chart_slide(
        prs, "PERFORMANCE COMPARISON", "PerformanceChart",
        XL_CHART_TYPE.BAR_CLUSTERED,
        categories=["Jan", "Feb", "Mar", "Apr", "May", "Jun"],
        series_data=[
            {"name": "Fund",      "values": (2.1, 0.9, 1.5, -0.3, 1.8, 2.4)},
            {"name": "Benchmark", "values": (1.8, 1.1, 1.3,  0.2, 1.5, 2.0)},
        ],
    )

    _make_chart_slide(
        prs, "CUMULATIVE RETURNS", "TrendChart",
        XL_CHART_TYPE.LINE,
        categories=["Q1", "Q2", "Q3", "Q4"],
        series_data=[
            {"name": "Fund",  "values": (3.2, 5.8, 7.1,  9.4)},
            {"name": "Index", "values": (2.8, 4.9, 6.3,  8.1)},
        ],
    )

    _make_chart_slide(
        prs, "ATTRIBUTION ANALYSIS", "AttributionChart",
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        categories=["Security\nSelection", "Sector\nAllocation", "Currency", "Other", "Total"],
        series_data=[
            {"name": "Contribution (bps)", "values": (45, 22, -8, 3, 62)},
        ],
    )

    out = OUTPUT_DIR / "potomac-chart-slides.pptx"
    prs.save(str(out))
    print(f"  \u2713 {out.name}")


# =============================================================================
# Template 3: potomac-table-slides.pptx
# =============================================================================

def build_table_slides():
    """
    3 slides:
      1. Data table     → TitleText, DataTable,     CaptionText
      2. Holdings table → TitleText, HoldingsTable, CaptionText
      3. Scorecard      → TitleText, ScorecardTable, CaptionText
    """
    prs = new_prs()

    # ── Slide 1: General Data Table ───────────────────────────────────────────
    slide = chart_slide_setup(prs, "DATA TABLE")
    add_table(
        slide, "DataTable",
        headers=["Name", "Value", "Change", "Weight"],
        rows=[
            ["Item A", "$100.00", "+2.1%", "25.0%"],
            ["Item B", "$85.50",  "+1.8%", "20.0%"],
            ["Item C", "$72.30",  "-0.5%", "18.0%"],
            ["Item D", "$68.90",  "+3.2%", "15.0%"],
            ["Item E", "$55.40",  "+0.9%", "12.0%"],
            ["TOTAL",  "$382.10", "+1.7%", "100.0%"],
        ],
        col_widths=[Inches(3.0), Inches(2.0), Inches(2.0), Inches(2.0)],
    )
    add_caption(slide)

    # ── Slide 2: Holdings Table ───────────────────────────────────────────────
    slide = chart_slide_setup(prs, "TOP HOLDINGS")
    add_table(
        slide, "HoldingsTable",
        headers=["Security", "Ticker", "Weight", "Return", "Contribution"],
        rows=[
            ["Apple Inc.",      "AAPL",  "8.2%", "+15.3%", "+1.25%"],
            ["Microsoft Corp.", "MSFT",  "7.1%", "+12.8%", "+0.91%"],
            ["NVIDIA Corp.",    "NVDA",  "5.8%", "+28.4%", "+1.65%"],
            ["Amazon.com Inc.", "AMZN",  "4.9%", "+8.6%",  "+0.42%"],
            ["Alphabet Inc.",   "GOOGL", "4.2%", "+11.2%", "+0.47%"],
        ],
        col_widths=[Inches(2.8), Inches(1.2), Inches(1.5), Inches(1.75), Inches(1.75)],
    )
    add_caption(slide,
                placeholder="As of {{date}} | {{num_holdings}} total positions | {{strategy}} strategy")

    # ── Slide 3: Scorecard Table ──────────────────────────────────────────────
    slide = chart_slide_setup(prs, "RISK SCORECARD")
    add_table(
        slide, "ScorecardTable",
        headers=["Metric", "Status", "Current", "Threshold", "Commentary"],
        rows=[
            ["Portfolio VaR",  "\u25CF ON TRACK", "1.2%",   "2.0%",  "Within limit"],
            ["Concentration",  "\u25CF AT RISK",   "18.5%",  "15.0%", "Monitor closely"],
            ["Drawdown",       "\u25CF ON TRACK", "-4.2%",  "-10.0%","Acceptable"],
            ["Beta",           "\u25CF ON TRACK", "0.87",   "1.20",  "Low market exposure"],
            ["Liquidity",      "\u25CF BREACH",    "82.1%",  "90.0%", "Needs rebalancing"],
        ],
        col_widths=[Inches(2.0), Inches(1.5), Inches(1.5), Inches(1.5), Inches(2.5)],
    )
    add_caption(slide,
                placeholder="Monitoring period: {{period}} | Last updated: {{date}}")

    out = OUTPUT_DIR / "potomac-table-slides.pptx"
    prs.save(str(out))
    print(f"  \u2713 {out.name}")


# =============================================================================
# Template 4: potomac-fund-fact-sheet.pptx
# =============================================================================

def build_fund_fact_sheet():
    """
    1 slide — complete fund fact sheet.
    Named shapes:
      FundName, AsOfDate, FundDescription,
      InceptionDate, AUM, NAV, YTDReturn, OneYearReturn, ThreeYearReturn,
      FiveYearReturn, SinceInception,
      BenchmarkName, BenchmarkYTD, Benchmark1Y, Benchmark3Y, Benchmark5Y,
      HoldingsTable,        (top holdings: Security, Weight, Return)
      AllocationChart,      (sector/asset allocation chart)
      RiskMetricsTable,     (Sharpe, Sortino, Max DD, Beta, StdDev)
      FooterText
    """
    prs = new_prs()
    slide = blank_slide(prs)
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE

    # ── Left yellow accent bar ────────────────────────────────────────────────
    rect(slide, "AccentBar",
         Inches(0), Inches(0), Inches(0.15), Inches(7.5), fill=YELLOW)

    # ── Header: Logo + Fund name ──────────────────────────────────────────────
    txbox(slide, "LogoArea",
          Inches(8.3), Inches(0.1), Inches(1.55), Inches(0.55),
          text="POTOMAC", font=FONT_H, size=14, color=DARK_GRAY,
          bold=True, align=PP_ALIGN.CENTER)

    txbox(slide, "FundName",
          Inches(0.35), Inches(0.1), Inches(7.5), Inches(0.55),
          text="{{fund_name}}",
          font=FONT_H, size=22, color=DARK_GRAY, bold=True)

    txbox(slide, "AsOfDate",
          Inches(0.35), Inches(0.65), Inches(4), Inches(0.35),
          text="As of {{date}}  |  Fact Sheet",
          font=FONT_B, size=12, color=GRAY_60)

    # Divider line
    rect(slide, "HeaderDivider",
         Inches(0.35), Inches(1.0), Inches(9.5), Inches(0.05), fill=YELLOW)

    # ── Description ───────────────────────────────────────────────────────────
    txbox(slide, "FundDescription",
          Inches(0.35), Inches(1.1), Inches(5.8), Inches(0.7),
          text="{{fund_description}}",
          font=FONT_B, size=10, color=DARK_GRAY)

    # ── Fund info boxes (right column) ────────────────────────────────────────
    info_items = [
        ("InceptionDate",  "Inception Date",  "{{inception_date}}"),
        ("AUM",            "AUM",             "{{aum}}"),
        ("NAV",            "NAV",             "{{nav}}"),
    ]
    for i, (name, label, val) in enumerate(info_items):
        y_pos = Inches(1.1 + i * 0.36)
        txbox(slide, f"InfoLabel_{name}",
              Inches(6.25), y_pos, Inches(1.5), Inches(0.3),
              text=label, font=FONT_B, size=9, color=GRAY_60, bold=True)
        txbox(slide, name,
              Inches(7.85), y_pos, Inches(1.95), Inches(0.3),
              text=val, font=FONT_B, size=9, color=DARK_GRAY, bold=True,
              align=PP_ALIGN.RIGHT)

    # ── Performance table ─────────────────────────────────────────────────────
    rect(slide, "PerfSectionBg",
         Inches(0.35), Inches(2.05), Inches(5.8), Inches(0.32), fill=YELLOW)
    txbox(slide, "PerfSectionLabel",
          Inches(0.35), Inches(2.05), Inches(5.8), Inches(0.32),
          text="PERFORMANCE", font=FONT_H, size=11,
          color=DARK_GRAY, bold=True, align=PP_ALIGN.CENTER)

    add_table(
        slide, "PerformanceTable",
        headers=["", "YTD", "1 Year", "3 Year", "5 Year", "Since Inception"],
        rows=[
            ["Fund",      "{{ytd}}", "{{1y}}", "{{3y}}", "{{5y}}", "{{si}}"],
            ["Benchmark", "{{bmk_ytd}}", "{{bmk_1y}}", "{{bmk_3y}}", "{{bmk_5y}}", "{{bmk_si}}"],
        ],
        col_widths=[Inches(1.3), Inches(0.9), Inches(0.9), Inches(0.9), Inches(0.9), Inches(0.9)],
        x=Inches(0.35), y=Inches(2.37),
    )

    txbox(slide, "BenchmarkName",
          Inches(0.35), Inches(3.32), Inches(5.8), Inches(0.3),
          text="Benchmark: {{benchmark_name}}",
          font=FONT_B, size=8, color=GRAY_60)

    # ── Holdings table ────────────────────────────────────────────────────────
    rect(slide, "HoldingsSectionBg",
         Inches(0.35), Inches(3.7), Inches(5.8), Inches(0.32), fill=YELLOW)
    txbox(slide, "HoldingsSectionLabel",
          Inches(0.35), Inches(3.7), Inches(5.8), Inches(0.32),
          text="TOP HOLDINGS", font=FONT_H, size=11,
          color=DARK_GRAY, bold=True, align=PP_ALIGN.CENTER)

    add_table(
        slide, "HoldingsTable",
        headers=["Security", "Weight", "Return"],
        rows=[
            ["{{holding_1_name}}", "{{holding_1_weight}}", "{{holding_1_return}}"],
            ["{{holding_2_name}}", "{{holding_2_weight}}", "{{holding_2_return}}"],
            ["{{holding_3_name}}", "{{holding_3_weight}}", "{{holding_3_return}}"],
            ["{{holding_4_name}}", "{{holding_4_weight}}", "{{holding_4_return}}"],
            ["{{holding_5_name}}", "{{holding_5_weight}}", "{{holding_5_return}}"],
        ],
        col_widths=[Inches(3.3), Inches(1.25), Inches(1.25)],
        x=Inches(0.35), y=Inches(4.02),
    )

    # ── Risk metrics (right side) ─────────────────────────────────────────────
    rect(slide, "RiskSectionBg",
         Inches(6.25), Inches(2.05), Inches(3.6), Inches(0.32), fill=YELLOW)
    txbox(slide, "RiskSectionLabel",
          Inches(6.25), Inches(2.05), Inches(3.6), Inches(0.32),
          text="RISK METRICS", font=FONT_H, size=11,
          color=DARK_GRAY, bold=True, align=PP_ALIGN.CENTER)

    add_table(
        slide, "RiskMetricsTable",
        headers=["Metric", "Value"],
        rows=[
            ["Sharpe Ratio",    "{{sharpe}}"],
            ["Sortino Ratio",   "{{sortino}}"],
            ["Max Drawdown",    "{{max_dd}}"],
            ["Beta",            "{{beta}}"],
            ["Std Deviation",   "{{std_dev}}"],
        ],
        col_widths=[Inches(2.0), Inches(1.6)],
        x=Inches(6.25), y=Inches(2.37),
    )

    # ── Allocation chart (right side, lower) ─────────────────────────────────
    rect(slide, "AllocSectionBg",
         Inches(6.25), Inches(3.7), Inches(3.6), Inches(0.32), fill=YELLOW)
    txbox(slide, "AllocSectionLabel",
          Inches(6.25), Inches(3.7), Inches(3.6), Inches(0.32),
          text="SECTOR ALLOCATION", font=FONT_H, size=11,
          color=DARK_GRAY, bold=True, align=PP_ALIGN.CENTER)

    # Allocation chart
    chart_data = CategoryChartData()
    chart_data.categories = ["Sector A", "Sector B", "Sector C", "Other"]
    chart_data.add_series("Allocation", (35.0, 28.0, 22.0, 15.0))

    gf = slide.shapes.add_chart(
        XL_CHART_TYPE.PIE,
        Inches(6.25), Inches(4.02),
        Inches(3.6), Inches(2.8),
        chart_data,
    )
    gf.name = "AllocationChart"
    gf.chart.has_legend = True

    # ── Footer ────────────────────────────────────────────────────────────────
    rect(slide, "FooterBar",
         Inches(0.35), Inches(7.1), Inches(9.5), Inches(0.05), fill=YELLOW)
    txbox(slide, "FooterText",
          Inches(0.35), Inches(7.18), Inches(9.5), Inches(0.3),
          text="Potomac Fund Management |  | For Advisor Use Only | "
               "Past performance is not indicative of future results.",
          font=FONT_B, size=7, color=GRAY_60, align=PP_ALIGN.CENTER)

    out = OUTPUT_DIR / "potomac-fund-fact-sheet.pptx"
    prs.save(str(out))
    print(f"  \u2713 {out.name}")


# =============================================================================
# Main
# =============================================================================

if __name__ == "__main__":
    print(f"\nGenerating Potomac automizer template files → {OUTPUT_DIR}\n")

    steps = [
        ("potomac-content-slides.pptx",  build_content_slides),
        ("potomac-chart-slides.pptx",    build_chart_slides),
        ("potomac-table-slides.pptx",    build_table_slides),
        ("potomac-fund-fact-sheet.pptx", build_fund_fact_sheet),
    ]

    for fname, fn in steps:
        try:
            fn()
        except Exception as exc:
            import traceback
            print(f"  \u2717 {fname} FAILED: {exc}")
            traceback.print_exc()

    print("\nShape name reference:")
    print("  content-slides :")
    print("    Slide 1: TitleText, SubtitleText, DateText, TaglineText")
    print("    Slide 2: TitleText, BodyText")
    print("    Slide 3: TitleText, LeftHeader, RightHeader, LeftContent, RightContent")
    print("    Slide 4: TitleText, Metric1-3 Value+Label, ContextText")
    print("    Slide 5: TitleText, DescriptionText")
    print("  chart-slides   :")
    print("    Slide 1: TitleText, PerformanceChart, CaptionText")
    print("    Slide 2: TitleText, TrendChart,       CaptionText")
    print("    Slide 3: TitleText, AttributionChart,  CaptionText")
    print("  table-slides   :")
    print("    Slide 1: TitleText, DataTable,      CaptionText")
    print("    Slide 2: TitleText, HoldingsTable,  CaptionText")
    print("    Slide 3: TitleText, ScorecardTable, CaptionText")
    print("  fund-fact-sheet:")
    print("    FundName, AsOfDate, FundDescription, InceptionDate, AUM, NAV,")
    print("    PerformanceTable, BenchmarkName, HoldingsTable,")
    print("    RiskMetricsTable, AllocationChart, FooterText")
    print("\nDone!")
