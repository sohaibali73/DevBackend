"""
File Preview API
================
Extracts structured preview data from PPTX and DOCX files stored in the
in-memory file_store (or fetched from Claude Files API).

Endpoints:
    GET /files/{file_id}/preview  — Returns slide/section data as JSON
"""

import io
import logging
from fastapi import APIRouter, HTTPException
from fastapi.responses import JSONResponse

from core.file_store import get_file

logger = logging.getLogger(__name__)

router = APIRouter(prefix="/files", tags=["preview"])


def _extract_pptx_slides(data: bytes) -> list:
    """Extract slide titles, text, and notes from a PPTX file."""
    try:
        from pptx import Presentation
    except ImportError:
        raise HTTPException(
            status_code=500,
            detail="python-pptx is not installed. Run: pip install python-pptx",
        )

    prs = Presentation(io.BytesIO(data))
    slides = []

    for idx, slide in enumerate(prs.slides, start=1):
        title = ""
        body_parts = []

        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text_frame.text.strip()
                if not text:
                    continue
                # First text shape that looks like a title
                if not title and (shape.shape_type == 13 or idx == 1 or len(text) < 120):
                    title = text
                else:
                    body_parts.append(text)

            # Tables
            if shape.has_table:
                table = shape.table
                rows = []
                for row in table.rows:
                    cells = [cell.text.strip() for cell in row.cells]
                    rows.append(" | ".join(cells))
                if rows:
                    body_parts.append("\n".join(rows))

        # Notes
        notes = ""
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            notes = slide.notes_slide.notes_text_frame.text.strip()

        slides.append({
            "slide_number": idx,
            "title": title or f"Slide {idx}",
            "content": "\n\n".join(body_parts) if body_parts else "",
            "notes": notes,
        })

    return slides


def _extract_docx_sections(data: bytes) -> list:
    """Extract sections/paragraphs from a DOCX file."""
    try:
        from docx import Document
    except ImportError:
        raise HTTPException(
            status_code=500,
            detail="python-docx is not installed. Run: pip install python-docx",
        )

    doc = Document(io.BytesIO(data))
    sections = []
    current_section = {"title": "", "content": ""}

    for para in doc.paragraphs:
        text = para.text.strip()
        if not text:
            continue

        # Detect headings
        if para.style and para.style.name and para.style.name.startswith("Heading"):
            # Save previous section if it has content
            if current_section["content"] or current_section["title"]:
                sections.append(current_section)
            current_section = {"title": text, "content": ""}
        else:
            if current_section["content"]:
                current_section["content"] += "\n" + text
            else:
                current_section["content"] = text

    # Don't forget the last section
    if current_section["content"] or current_section["title"]:
        sections.append(current_section)

    return sections


def _try_claude_files_fallback(file_id: str) -> tuple:
    """Try to download a file from Claude's Files API. Returns (data, filename) or (None, None)."""
    if not file_id.startswith("file_"):
        return None, None

    try:
        import anthropic
        import os

        api_key = os.getenv("ANTHROPIC_API_KEY")
        if not api_key:
            try:
                from config import settings
                api_key = getattr(settings, "ANTHROPIC_API_KEY", None) or getattr(settings, "anthropic_api_key", None)
            except Exception:
                pass

        if not api_key:
            return None, None

        client = anthropic.Anthropic(api_key=api_key)
        files_beta = ["files-api-2025-04-14"]

        metadata = client.beta.files.retrieve_metadata(file_id, betas=files_beta)
        filename = getattr(metadata, "filename", None) or f"download_{file_id}"

        file_content = client.beta.files.download(file_id, betas=files_beta)
        if isinstance(file_content, bytes):
            data = file_content
        elif hasattr(file_content, "iter_bytes"):
            data = b"".join(file_content.iter_bytes())
        elif hasattr(file_content, "read"):
            data = file_content.read()
        else:
            data = bytes(file_content)

        # Cache for future requests
        from core.file_store import store_file
        store_file(data=data, filename=filename, tool_name="claude_files_api", file_id=file_id)

        logger.info("Downloaded file %s from Claude API for preview (%d bytes)", file_id, len(data))
        return data, filename

    except Exception as e:
        logger.warning("Claude Files API fallback failed for %s: %s", file_id, e)
        return None, None


@router.get("/{file_id}/preview")
async def preview_file(file_id: str):
    """
    Extract structured preview data from a PPTX or DOCX file.

    Returns:
        For PPTX: { "type": "pptx", "slides": [...], "slide_count": N }
        For DOCX: { "type": "docx", "sections": [...], "page_count": N }
    """
    entry = get_file(file_id)

    # If not in memory, try Claude Files API fallback
    if entry is None:
        data, filename = _try_claude_files_fallback(file_id)
        if data is not None and filename is not None:
            # Create a simple namespace to mimic file entry
            class _Entry:
                pass
            entry = _Entry()
            entry.data = data
            entry.filename = filename

    if entry is None:
        raise HTTPException(status_code=404, detail=f"File {file_id} not found")

    filename = entry.filename.lower()
    data = entry.data

    if filename.endswith(".pptx") or filename.endswith(".ppt"):
        try:
            slides = _extract_pptx_slides(data)
            return JSONResponse({
                "type": "pptx",
                "filename": entry.filename,
                "slides": slides,
                "slide_count": len(slides),
            })
        except HTTPException:
            raise
        except Exception as e:
            logger.error("PPTX preview extraction failed: %s", e)
            raise HTTPException(status_code=500, detail=f"Failed to extract PPTX preview: {e}")

    elif filename.endswith(".docx") or filename.endswith(".doc"):
        try:
            sections = _extract_docx_sections(data)
            return JSONResponse({
                "type": "docx",
                "filename": entry.filename,
                "sections": sections,
                "page_count": len(sections),
            })
        except HTTPException:
            raise
        except Exception as e:
            logger.error("DOCX preview extraction failed: %s", e)
            raise HTTPException(status_code=500, detail=f"Failed to extract DOCX preview: {e}")

    else:
        raise HTTPException(
            status_code=400,
            detail=f"Preview not supported for file type: {entry.filename}",
        )
